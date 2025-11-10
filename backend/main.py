"""
AI Study Assistant Backend - FastAPI Application
Provides grounded document processing, exam generation, and AI tutoring
"""
from fastapi import FastAPI, Depends, HTTPException, UploadFile, File, Request, Header, Response
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from sqlalchemy import create_engine, Column, Integer, String, DateTime, Float, Date, text, func
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, Session
from datetime import datetime, timedelta, date
from typing import List, Optional, Dict, Any
from pydantic import BaseModel, EmailStr
import bcrypt
from jose import jwt, JWTError
import os
import re
import requests
from collections import defaultdict
import stripe
import time
from dotenv import load_dotenv
import io
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

# Load environment variables from .env file
load_dotenv()

# ============================================================================
# CONFIGURATION
# ============================================================================
DATABASE_URL = os.getenv("DATABASE_URL", "sqlite:///./study_assistant.db")
SECRET_KEY = os.getenv("JWT_SECRET_KEY", "your-secret-key-change-in-production")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 43200  # 30 days (30 * 24 * 60)
REFRESH_TOKEN_EXPIRE_DAYS = 90  # 90 days
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
STRIPE_SECRET_KEY = os.getenv("STRIPE_SECRET_KEY")
STRIPE_WEBHOOK_SECRET = os.getenv("STRIPE_WEBHOOK_SECRET")

# Cookie settings - Railway/production detection
# Railway always sets PORT env var, use that as production indicator
IS_PRODUCTION = bool(os.getenv("RAILWAY_ENVIRONMENT") or os.getenv("RENDER") or os.getenv("PORT"))

# For Railway (production), always use secure cookies with SameSite=none for cross-origin (Vercel + Railway)
# Railway provides HTTPS by default, so secure=True is safe
COOKIE_SECURE = IS_PRODUCTION  # True in production (Railway), False in local dev
COOKIE_SAMESITE = "none" if IS_PRODUCTION else "lax"

print(f"[AUTH CONFIG] IS_PRODUCTION={IS_PRODUCTION}, COOKIE_SECURE={COOKIE_SECURE}, COOKIE_SAMESITE={COOKIE_SAMESITE}")

if STRIPE_SECRET_KEY:
    stripe.api_key = STRIPE_SECRET_KEY

# Rate limiting storage (in-memory, use Redis in production)
rate_limit_store = defaultdict(list)

# File content storage for anonymous users (in-memory)
file_content_store = {}

# ============================================================================
# DATABASE SETUP
# ============================================================================
# check_same_thread is only for SQLite
try:
    if DATABASE_URL.startswith("sqlite"):
        engine = create_engine(DATABASE_URL, connect_args={"check_same_thread": False})
    else:
        # Optimized for long-running operations and Supabase connection limits
        engine = create_engine(
            DATABASE_URL,
            pool_size=3,              # Reduced for Supabase Session mode limits
            max_overflow=5,           # Total max = 8 connections
            pool_recycle=300,         # Recycle connections after 5 minutes
            pool_pre_ping=True,       # Check connection health before use
            pool_timeout=10,          # Timeout waiting for connection (seconds)
            connect_args={
                "connect_timeout": 10,  # PostgreSQL connection timeout
                "keepalives": 1,        # Enable TCP keepalives
                "keepalives_idle": 30,  # Start keepalives after 30s idle
                "keepalives_interval": 10,  # Keepalive interval
                "keepalives_count": 5   # Max keepalive probes
            }
        )
    SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
    Base = declarative_base()
    print(f"[DATABASE] Database engine created successfully (URL: {DATABASE_URL[:20]}...)")
except Exception as e:
    print(f"[DATABASE ERROR] Failed to create database engine: {e}")
    import traceback
    traceback.print_exc()
    # Create a dummy engine to prevent crashes - migrations will handle actual connection
    engine = None
    SessionLocal = None
    Base = declarative_base()
    print("[DATABASE WARNING] Using fallback database setup - some features may not work")

class User(Base):
    __tablename__ = "users"
    id = Column(Integer, primary_key=True, index=True)
    email = Column(String, unique=True, index=True)
    password_hash = Column(String)
    name = Column(String)
    surname = Column(String)
    oauth_provider = Column(String, nullable=True)  # Legacy - not used
    oauth_id = Column(String, nullable=True)  # Legacy - not used
    tier = Column(String, default="free")  # "free" or "premium"
    is_admin = Column(Integer, default=0)  # 0 = false, 1 = true (SQLite compatibility)
    created_at = Column(DateTime, default=datetime.utcnow)

class Upload(Base):
    __tablename__ = "uploads"
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer)
    file_id = Column(String)  # OpenAI file ID (not used anymore)
    filename = Column(String)
    mime = Column(String)
    size = Column(Integer)
    content = Column(String)  # Extracted text content
    created_at = Column(DateTime, default=datetime.utcnow)

class Usage(Base):
    __tablename__ = "usage"
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer)
    kind = Column(String)  # "exam", "explain", "chat", "upload"
    count = Column(Integer, default=1)
    date = Column(Date, default=date.today)

class Exam(Base):
    __tablename__ = "exams"
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer)
    payload_json = Column(String)
    created_at = Column(DateTime, default=datetime.utcnow)

class Folder(Base):
    __tablename__ = "folders"
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer)
    name = Column(String)
    color = Column(String, nullable=True)  # Hex color for folder (optional)
    icon = Column(String, nullable=True)  # Emoji or icon name (optional)
    created_at = Column(DateTime, default=datetime.utcnow)

class History(Base):
    __tablename__ = "history"
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer)
    folder_id = Column(Integer, nullable=True)  # Reference to folders table
    type = Column(String)  # "summary", "flashcards", "truefalse", "exam"
    title = Column(String)
    data_json = Column(String)  # JSON stringified data
    score_json = Column(String, nullable=True)  # JSON stringified score (for truefalse, exam)
    created_at = Column(DateTime, default=datetime.utcnow)

class Transaction(Base):
    __tablename__ = "transactions"
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, index=True)
    stripe_session_id = Column(String, unique=True, index=True)
    stripe_customer_id = Column(String, nullable=True)
    stripe_subscription_id = Column(String, nullable=True)
    stripe_payment_intent_id = Column(String, nullable=True)
    amount = Column(Float)  # Amount in cents (Stripe format)
    currency = Column(String, default="usd")
    status = Column(String)  # "completed", "pending", "failed", "refunded"
    tier = Column(String)  # "premium", "pro", etc.
    event_type = Column(String)  # "checkout.session.completed", "invoice.paid", etc.
    event_metadata = Column(String, nullable=True)  # JSON string for additional data (renamed from 'metadata' because it's reserved in SQLAlchemy)
    created_at = Column(DateTime, default=datetime.utcnow)
    updated_at = Column(DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)

class TokenUsage(Base):
    __tablename__ = "token_usage"
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, nullable=True, index=True)  # Null for anonymous users
    endpoint = Column(String)  # "summarize", "flashcards", "exam", "chat", "explain"
    model = Column(String, default="gpt-4o-mini")
    input_tokens = Column(Integer, default=0)
    output_tokens = Column(Integer, default=0)
    total_tokens = Column(Integer, default=0)
    estimated_cost = Column(Float, nullable=True)  # Estimated cost in USD
    created_at = Column(DateTime, default=datetime.utcnow)

# Create tables if engine is available (non-blocking)
try:
    if engine:
        Base.metadata.create_all(bind=engine)
        print("[DATABASE] Tables created/verified successfully")
    else:
        print("[DATABASE WARNING] Skipping table creation - engine not available")
except Exception as e:
    print(f"[DATABASE WARNING] Table creation failed: {e}")
    # Don't crash - tables might already exist or connection will be retried

# ============================================================================
# DATABASE MIGRATION FUNCTION
# ============================================================================
def run_migration():
    """Run database migrations to ensure all required columns and tables exist"""
    if not engine:
        print("[MIGRATION] Skipping migration - database engine not available")
        return
    
    try:
        from sqlalchemy import inspect
        inspector = inspect(engine)
        existing_tables = []
        users_columns = []
        history_columns = []
        
        try:
            existing_tables = inspector.get_table_names()
            
            # Check if users table exists and get its columns
            if "users" in existing_tables:
                try:
                    users_columns = [col["name"].lower() for col in inspector.get_columns("users")]
                except Exception as e:
                    print(f"[MIGRATION] Could not get users table columns: {e}")
                    users_columns = []
            
            # Check history table columns
            if "history" in existing_tables:
                try:
                    history_columns = [col["name"].lower() for col in inspector.get_columns("history")]
                except Exception as e:
                    print(f"[MIGRATION] Could not get history table columns: {e}")
                    history_columns = []
        except Exception as e:
            print(f"[MIGRATION] Could not inspect database: {e}")
            # Continue with migration anyway - will handle errors gracefully
        
        # Create a database session for migration
        db = SessionLocal()
        try:
            # Import telemetry model to register it
            try:
                from app.models.telemetry import SummaryQuality
                # Create telemetry table if missing
                if "summary_quality" not in existing_tables:
                    SummaryQuality.__table__.create(engine)
                    print("[MIGRATION] Created summary_quality table")
            except Exception as e:
                print(f"[MIGRATION] Could not import telemetry model: {e}")
            
            # Add missing columns to users table
            if DATABASE_URL.startswith("sqlite"):
                # SQLite
                if "name" not in users_columns:
                    try:
                        db.execute(text("ALTER TABLE users ADD COLUMN name TEXT"))
                        db.commit()
                        print("[MIGRATION] Added 'name' column to users table")
                    except Exception as e:
                        db.rollback()
                        print(f"[MIGRATION] Could not add 'name' column: {e}")
                
                if "surname" not in users_columns:
                    try:
                        db.execute(text("ALTER TABLE users ADD COLUMN surname TEXT"))
                        db.commit()
                        print("[MIGRATION] Added 'surname' column to users table")
                    except Exception as e:
                        db.rollback()
                        print(f"[MIGRATION] Could not add 'surname' column: {e}")
                
                if "is_admin" not in users_columns:
                    try:
                        db.execute(text("ALTER TABLE users ADD COLUMN is_admin INTEGER DEFAULT 0"))
                        db.commit()
                        print("[MIGRATION] Added 'is_admin' column to users table")
                    except Exception as e:
                        db.rollback()
                        # Check if column already exists (might have been added between check and alter)
                        error_str = str(e).lower()
                        if "duplicate column" in error_str or "already exists" in error_str:
                            print("[MIGRATION] 'is_admin' column already exists")
                        else:
                            print(f"[MIGRATION] Could not add 'is_admin' column: {e}")
                
                if "folder_id" not in history_columns:
                    try:
                        db.execute(text("ALTER TABLE history ADD COLUMN folder_id INTEGER"))
                        db.commit()
                        print("[MIGRATION] Added 'folder_id' column to history table")
                    except Exception as e:
                        db.rollback()
                        print(f"[MIGRATION] Could not add 'folder_id' column: {e}")
                
                # Create folders table
                if "folders" not in existing_tables:
                    try:
                        db.execute(text("""
                            CREATE TABLE IF NOT EXISTS folders (
                                id INTEGER PRIMARY KEY,
                                user_id INTEGER,
                                name TEXT,
                                color TEXT,
                                icon TEXT,
                                created_at TIMESTAMP
                            )
                        """))
                        db.commit()
                        print("[MIGRATION] Created 'folders' table")
                    except Exception as e:
                        db.rollback()
                        print(f"[MIGRATION] Could not create 'folders' table: {e}")
                
                # Create transactions table
                if "transactions" not in existing_tables:
                    try:
                        db.execute(text("""
                            CREATE TABLE IF NOT EXISTS transactions (
                                id INTEGER PRIMARY KEY,
                                user_id INTEGER,
                                stripe_session_id TEXT UNIQUE,
                                stripe_customer_id TEXT,
                                stripe_subscription_id TEXT,
                                stripe_payment_intent_id TEXT,
                                amount REAL,
                                currency TEXT DEFAULT 'usd',
                                status TEXT,
                                tier TEXT,
                                event_type TEXT,
                                event_metadata TEXT,
                                created_at TIMESTAMP,
                                updated_at TIMESTAMP
                            )
                        """))
                        db.commit()
                        print("[MIGRATION] Created 'transactions' table")
                    except Exception as e:
                        db.rollback()
                        print(f"[MIGRATION] Could not create 'transactions' table: {e}")
                
                # Create token_usage table
                if "token_usage" not in existing_tables:
                    try:
                        db.execute(text("""
                            CREATE TABLE IF NOT EXISTS token_usage (
                                id INTEGER PRIMARY KEY,
                                user_id INTEGER,
                                endpoint TEXT,
                                model TEXT DEFAULT 'gpt-4o-mini',
                                input_tokens INTEGER DEFAULT 0,
                                output_tokens INTEGER DEFAULT 0,
                                total_tokens INTEGER DEFAULT 0,
                                estimated_cost REAL,
                                created_at TIMESTAMP
                            )
                        """))
                        db.commit()
                        print("[MIGRATION] Created 'token_usage' table")
                    except Exception as e:
                        db.rollback()
                        print(f"[MIGRATION] Could not create 'token_usage' table: {e}")
            else:
                # PostgreSQL
                try:
                    if "name" not in users_columns:
                        db.execute(text("ALTER TABLE users ADD COLUMN IF NOT EXISTS name VARCHAR"))
                        db.commit()
                        print("[MIGRATION] Added 'name' column to users table")
                    
                    if "surname" not in users_columns:
                        db.execute(text("ALTER TABLE users ADD COLUMN IF NOT EXISTS surname VARCHAR"))
                        db.commit()
                        print("[MIGRATION] Added 'surname' column to users table")
                    
                    if "is_admin" not in users_columns:
                        try:
                            db.execute(text("ALTER TABLE users ADD COLUMN IF NOT EXISTS is_admin INTEGER DEFAULT 0"))
                            db.commit()
                            print("[MIGRATION] Added 'is_admin' column to users table")
                        except Exception as e:
                            db.rollback()
                            print(f"[MIGRATION] Could not add 'is_admin' column: {e}")
                    
                    if "folder_id" not in history_columns:
                        db.execute(text("ALTER TABLE history ADD COLUMN IF NOT EXISTS folder_id INTEGER"))
                        db.commit()
                        print("[MIGRATION] Added 'folder_id' column to history table")
                    
                    # Create folders table
                    if "folders" not in existing_tables:
                        db.execute(text("""
                            CREATE TABLE IF NOT EXISTS folders (
                                id SERIAL PRIMARY KEY,
                                user_id INTEGER,
                                name VARCHAR,
                                color VARCHAR,
                                icon VARCHAR,
                                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                            )
                        """))
                        db.commit()
                        print("[MIGRATION] Created 'folders' table")
                    
                    # Create transactions table
                    if "transactions" not in existing_tables:
                        db.execute(text("""
                            CREATE TABLE IF NOT EXISTS transactions (
                                id SERIAL PRIMARY KEY,
                                user_id INTEGER,
                                stripe_session_id VARCHAR UNIQUE,
                                stripe_customer_id VARCHAR,
                                stripe_subscription_id VARCHAR,
                                stripe_payment_intent_id VARCHAR,
                                amount REAL,
                                currency VARCHAR DEFAULT 'usd',
                                status VARCHAR,
                                tier VARCHAR,
                                event_type VARCHAR,
                                event_metadata TEXT,
                                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                                updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                            )
                        """))
                        db.commit()
                        print("[MIGRATION] Created 'transactions' table")
                    
                    # Create token_usage table
                    if "token_usage" not in existing_tables:
                        db.execute(text("""
                            CREATE TABLE IF NOT EXISTS token_usage (
                                id SERIAL PRIMARY KEY,
                                user_id INTEGER,
                                endpoint VARCHAR,
                                model VARCHAR DEFAULT 'gpt-4o-mini',
                                input_tokens INTEGER DEFAULT 0,
                                output_tokens INTEGER DEFAULT 0,
                                total_tokens INTEGER DEFAULT 0,
                                estimated_cost REAL,
                                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                            )
                        """))
                        db.commit()
                        print("[MIGRATION] Created 'token_usage' table")
                except Exception as e:
                    db.rollback()
                    print(f"[MIGRATION] Error during PostgreSQL migration: {e}")
            
            print("[MIGRATION] Database migration completed successfully")
        finally:
            db.close()
    except Exception as e:
        print(f"[MIGRATION] Migration error: {e}")
        import traceback
        traceback.print_exc()

# Run migration on startup (non-blocking - app will start even if migration fails)
print("[STARTUP] Running database migration...")
try:
    run_migration()
    print("[STARTUP] Migration completed successfully")
except Exception as e:
    print(f"[STARTUP WARNING] Migration failed but continuing startup: {e}")
    import traceback
    traceback.print_exc()
    # Don't crash the app - migration can be run manually later

# ============================================================================
# FASTAPI APP
# ============================================================================
app = FastAPI(title="AI Study Assistant API", version="1.0.0")

# Get CORS origins from environment variable
cors_origins_str = os.getenv("CORS_ORIGINS", "")
if cors_origins_str:
    cors_origins = [origin.strip() for origin in cors_origins_str.split(",") if origin.strip()]
else:
    # Default to allow all origins if not specified
    cors_origins = ["*"]

# In production environments, always allow all origins for now
# TODO: Restrict to specific domains in production
if os.getenv("RAILWAY_ENVIRONMENT") or os.getenv("RENDER") or os.getenv("PORT"):
    cors_origins = ["*"]

# Add CORS middleware with permissive settings
app.add_middleware(
    CORSMiddleware,
    allow_origins=cors_origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["*"],
    max_age=3600,  # Cache preflight requests for 1 hour
)

# Add custom CORS headers to all responses as backup
@app.middleware("http")
async def add_cors_headers(request, call_next):
    # Get origin from request
    origin = request.headers.get("origin")
    
    # Handle preflight requests
    if request.method == "OPTIONS":
        response = JSONResponse(content={}, status_code=200)
        if origin:
            response.headers["Access-Control-Allow-Origin"] = origin
        response.headers["Access-Control-Allow-Credentials"] = "true"
        response.headers["Access-Control-Allow-Methods"] = "GET, POST, PUT, DELETE, OPTIONS, PATCH"
        response.headers["Access-Control-Allow-Headers"] = "Content-Type, Authorization, Accept, Origin, X-Requested-With"
        response.headers["Access-Control-Max-Age"] = "3600"
        return response
    
    # Let FastAPI handle the request and catch HTTPExceptions properly
    response = await call_next(request)
    
    # Set specific origin instead of wildcard when credentials are used
    if origin:
        response.headers["Access-Control-Allow-Origin"] = origin
    response.headers["Access-Control-Allow-Credentials"] = "true"
    response.headers["Access-Control-Allow-Methods"] = "GET, POST, PUT, DELETE, OPTIONS, PATCH"
    response.headers["Access-Control-Allow-Headers"] = "Content-Type, Authorization, Accept, Origin, X-Requested-With"
    response.headers["Access-Control-Expose-Headers"] = "Content-Type, Authorization"
    return response

# ============================================================================
# PYDANTIC MODELS
# ============================================================================
class UserRegister(BaseModel):
    email: EmailStr
    password: str
    name: str
    surname: str

class UserLogin(BaseModel):
    email: EmailStr
    password: str

class TokenResponse(BaseModel):
    access_token: str
    refresh_token: str
    token_type: str = "bearer"

class SummarizeRequest(BaseModel):
    file_ids: Optional[List[str]] = None
    language: Optional[str] = "en"
    outline: Optional[bool] = False
    prompt: Optional[str] = None

class FlashcardsRequest(BaseModel):
    file_ids: Optional[List[str]] = None
    style: Optional[str] = "basic"
    deck_name: Optional[str] = "Study Deck"
    count: Optional[int] = 10
    language: Optional[str] = "en"
    prompt: Optional[str] = None

class TrueFalseRequest(BaseModel):
    file_ids: Optional[List[str]] = None
    count: Optional[int] = 10
    language: Optional[str] = "en"
    prompt: Optional[str] = None

class ExamRequest(BaseModel):
    file_ids: Optional[List[str]] = None
    level: Optional[str] = "high-school"
    count: Optional[int] = 5
    language: Optional[str] = "en"
    prompt: Optional[str] = None

class AskRequest(BaseModel):
    prompt: str
    level: Optional[str] = "high-school"
    count: Optional[int] = 5

class ExplainRequest(BaseModel):
    question: str
    options: Optional[Dict[str, str]] = None
    selected: Optional[str] = None
    correct: Optional[str] = None
    file_ids: Optional[List[str]] = None

class ChatMessage(BaseModel):
    role: str
    content: str

class ChatRequest(BaseModel):
    messages: List[ChatMessage]
    file_ids: Optional[List[str]] = None

class CheckoutRequest(BaseModel):
    priceId: str
    successUrl: str
    cancelUrl: str

class FolderCreate(BaseModel):
    name: str
    color: Optional[str] = None
    icon: Optional[str] = None

class FolderUpdate(BaseModel):
    name: Optional[str] = None
    color: Optional[str] = None
    icon: Optional[str] = None

class FolderResponse(BaseModel):
    id: int
    name: str
    color: Optional[str] = None
    icon: Optional[str] = None
    created_at: str
    item_count: int = 0

class HistoryItemCreate(BaseModel):
    type: str  # "summary", "flashcards", "truefalse", "exam"
    title: str
    data: Any  # Will be JSON stringified
    score: Optional[Any] = None  # For truefalse and exam results
    folder_id: Optional[int] = None

class HistoryItemUpdate(BaseModel):
    title: Optional[str] = None
    data: Optional[Any] = None
    score: Optional[Any] = None
    folder_id: Optional[int] = None

class HistoryItemResponse(BaseModel):
    id: int
    type: str
    title: str
    data: Any
    score: Optional[Any] = None
    folder_id: Optional[int] = None
    timestamp: int  # Unix timestamp in milliseconds

class UserUpdate(BaseModel):
    name: Optional[str] = None
    surname: Optional[str] = None
    tier: Optional[str] = None
    is_admin: Optional[bool] = None

class BootstrapAdminRequest(BaseModel):
    email: str
    secret_key: Optional[str] = None

class UserResponse(BaseModel):
    id: int
    email: str
    name: Optional[str] = None
    surname: Optional[str] = None
    tier: str
    is_admin: bool
    created_at: str
    usage: Optional[Dict[str, int]] = None
    history_count: Optional[int] = None
    uploads_count: Optional[int] = None

# ============================================================================
# DEPENDENCIES
# ============================================================================
def get_db():
    if not SessionLocal:
        raise HTTPException(status_code=503, detail="Database connection not available. Please check backend configuration.")
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

def get_current_user(request: Request, authorization: Optional[str] = Header(None), db: Session = Depends(get_db)) -> User:
    # Try to get token from cookie first (preferred), then fall back to Authorization header
    token = request.cookies.get("access_token")
    
    # Fallback to Authorization header for backward compatibility
    if not token and authorization and authorization.startswith("Bearer "):
        token = authorization.split(" ")[1]
    
    if not token:
        raise HTTPException(status_code=401, detail="Missing or invalid authorization")
    
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        user_id_str = payload.get("sub")
        if user_id_str is None:
            raise HTTPException(status_code=401, detail="Invalid token")
        user_id = int(user_id_str)  # Convert string back to int
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail="Token expired")
    except JWTError:
        raise HTTPException(status_code=401, detail="Invalid token")
    
    try:
        # Query user - use explicit column selection to avoid issues
        user = db.query(User).filter(User.id == user_id).first()
        if not user:
            raise HTTPException(status_code=401, detail="User not found")
        return user
    except Exception as e:
        # Log the error for debugging
        import traceback
        print(f"[AUTH ERROR] Failed to get user {user_id}: {str(e)}")
        print(f"[AUTH ERROR] Traceback: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Failed to retrieve user information")

def get_optional_user(request: Request, authorization: Optional[str] = Header(None), db: Session = Depends(get_db)) -> Optional[User]:
    """Get user if authenticated, None otherwise"""
    # Try to get token from cookie first, then fall back to Authorization header
    token = request.cookies.get("access_token")
    
    if not token and authorization and authorization.startswith("Bearer "):
        token = authorization.split(" ")[1]
    
    if not token:
        return None
    
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        user_id_str = payload.get("sub")
        if user_id_str is None:
            return None
        
        user_id = int(user_id_str)  # Convert string back to int
        user = db.query(User).filter(User.id == user_id).first()
        return user
    except:
        return None

def get_admin_user(current_user: User = Depends(get_current_user)) -> User:
    """Get current user and verify they are an admin"""
    # Check if user is admin (is_admin can be 1 or True)
    if not current_user.is_admin or current_user.is_admin == 0:
        raise HTTPException(status_code=403, detail="Admin access required")
    return current_user

# ============================================================================
# UTILITIES
# ============================================================================
def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(password: str, hashed: str) -> bool:
    return bcrypt.checkpw(password.encode('utf-8'), hashed.encode('utf-8'))

def create_token(user_id: int, expires_delta: timedelta) -> str:
    expire = datetime.utcnow() + expires_delta
    payload = {"sub": str(user_id), "exp": expire}  # Convert user_id to string
    return jwt.encode(payload, SECRET_KEY, algorithm=ALGORITHM)

def rate_limit(request: Request, limit: int = 30, window: int = 300):
    """Simple in-memory rate limiting: limit requests per window (seconds)"""
    ip = request.client.host
    now = time.time()
    
    # Clean old entries
    rate_limit_store[ip] = [ts for ts in rate_limit_store[ip] if now - ts < window]
    
    if len(rate_limit_store[ip]) >= limit:
        raise HTTPException(status_code=429, detail="Rate limit exceeded")
    
    rate_limit_store[ip].append(now)

def check_quota(db: Session, user: User, kind: str) -> bool:
    """Check if user has quota remaining for the given action"""
    # QUOTAS DISABLED FOR TESTING - Always return True
    return True
    
    # Original quota logic (disabled for testing)
    # if user.tier == "premium":
    #     # Premium users have generous limits
    #     limits = {"exam": 100, "explain": 500, "chat": 1000, "upload": 100}
    # else:
    #     # Free tier limits
    #     limits = {"exam": 2, "explain": 5, "chat": 10, "upload": 2}
    # 
    # today = date.today()
    # usage = db.query(Usage).filter(
    #     Usage.user_id == user.id,
    #     Usage.kind == kind,
    #     Usage.date == today
    # ).first()
    # 
    # current_count = usage.count if usage else 0
    # return current_count < limits.get(kind, 0)

def increment_usage(db: Session, user: User, kind: str):
    """Increment usage counter for the user"""
    today = date.today()
    usage = db.query(Usage).filter(
        Usage.user_id == user.id,
        Usage.kind == kind,
        Usage.date == today
    ).first()
    
    if usage:
        usage.count += 1
    else:
        usage = Usage(user_id=user.id, kind=kind, count=1, date=today)
        db.add(usage)
    
    db.commit()

def get_level_text(level: str) -> str:
    """Map difficulty level to description text"""
    mapping = {
        "elementary-middle": "elementary/middle-school level (core notions, simple language)",
        "high-school": "high-school level (intermediate concepts and applications)",
        "university": "university level (advanced concepts; proofs/applications)"
    }
    return mapping.get(level, "university level (advanced concepts; proofs/applications)")

# ============================================================================
# OPENAI INTEGRATION
# ============================================================================
def extract_text_from_file(file_content: bytes, filename: str, mime: str) -> str:
    """Extract text content from uploaded files"""
    try:
        # PDF
        if mime == "application/pdf" or filename.endswith('.pdf'):
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PdfReader(pdf_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        
        # DOCX
        elif mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or filename.endswith('.docx'):
            docx_file = io.BytesIO(file_content)
            doc = Document(docx_file)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text.strip()
        
        # PPTX
        elif mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or filename.endswith('.pptx'):
            pptx_file = io.BytesIO(file_content)
            prs = Presentation(pptx_file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        
        # Plain text or unknown - try to decode as text
        else:
            try:
                return file_content.decode('utf-8')
            except:
                return file_content.decode('latin-1')
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to extract text from {filename}: {str(e)}")

def call_openai_with_context(file_contents: List[str], prompt: str, temperature: float = 0.0, model: str = "gpt-4o-mini", max_tokens: int = 4000, user_id: Optional[int] = None, endpoint: str = "unknown", db: Optional[Session] = None) -> str:
    """Call OpenAI API with file contents included in the prompt. Returns response text."""
    if not OPENAI_API_KEY:
        raise HTTPException(status_code=500, detail="OpenAI API key not configured")
    
    url = "https://api.openai.com/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {OPENAI_API_KEY}",
        "Content-Type": "application/json"
    }
    
    # Build messages with system and user roles
    messages = []
    
    if file_contents:
        # Add system message with strict instructions
        messages.append({
            "role": "system",
            "content": "You are a study assistant. You MUST create content based ONLY on the documents provided. Do not use external knowledge. If you generate exam questions, they MUST be about the specific content in the documents, not general topics."
        })
        
        # Add document context
        context = "DOCUMENT CONTENT:\n\n"
        for i, content in enumerate(file_contents, 1):
            context += f"=== Document {i} ===\n{content}\n\n"
        
        context += f"\n{prompt}\n\nREMEMBER: Base your response ONLY on the document content above. Do not use external knowledge."
        
        messages.append({
            "role": "user",
            "content": context
        })
    else:
        messages.append({
            "role": "user",
            "content": prompt
        })
    
    payload = {
        "model": model,
        "messages": messages,
        "temperature": temperature,
        "max_tokens": max_tokens
    }
    
    response = requests.post(url, headers=headers, json=payload, timeout=60)
    
    if response.status_code != 200:
        raise HTTPException(status_code=500, detail=f"OpenAI API call failed: {response.text}")
    
    response_data = response.json()
    content = response_data["choices"][0]["message"]["content"]
    usage = response_data.get("usage", {})
    
    # Track token usage in database (non-blocking)
    if db and endpoint != "unknown":
        try:
            input_tokens = usage.get("prompt_tokens", 0)
            output_tokens = usage.get("completion_tokens", 0)
            total_tokens = usage.get("total_tokens", 0)
            
            # Cost calculation (per 1M tokens)
            if "gpt-4o" in model.lower():
                input_cost_per_1m = 2.50  # $2.50 per 1M input tokens for gpt-4o
                output_cost_per_1m = 10.00  # $10.00 per 1M output tokens for gpt-4o
            elif "gpt-4" in model.lower():
                input_cost_per_1m = 30.00  # $30 per 1M input tokens for gpt-4
                output_cost_per_1m = 60.00  # $60 per 1M output tokens for gpt-4
            else:
                input_cost_per_1m = 0.150  # $0.15 per 1M input tokens for gpt-4o-mini
                output_cost_per_1m = 0.600  # $0.60 per 1M output tokens for gpt-4o-mini
            
            estimated_cost = (input_tokens / 1_000_000 * input_cost_per_1m) + (output_tokens / 1_000_000 * output_cost_per_1m)
            
            token_usage = TokenUsage(
                user_id=user_id,
                endpoint=endpoint,
                model=model,
                input_tokens=input_tokens,
                output_tokens=output_tokens,
                total_tokens=total_tokens,
                estimated_cost=estimated_cost
            )
            db.add(token_usage)
            db.commit()
        except Exception as e:
            # Don't fail the request if token tracking fails
            print(f"[TOKEN TRACKING ERROR] Failed to record token usage: {e}")
            try:
                db.rollback()
            except:
                pass
    
    return content

def parse_mcq_questions(text: str) -> tuple:
    """Parse MCQ format and return questions, answer key"""
    questions = []
    answer_key = {}
    
    # Split by answer key section
    parts = re.split(r'\n\s*Cevap Anahtarı:\s*\n', text, flags=re.IGNORECASE)
    
    if len(parts) == 2:
        questions_text, answers_text = parts
    else:
        questions_text = text
        answers_text = ""
    
    # Parse questions
    question_pattern = r'(\d+)\.\s+(.*?)(?=\n\d+\.|\nCevap Anahtarı:|$)'
    question_matches = re.findall(question_pattern, questions_text, re.DOTALL)
    
    for number, q_text in question_matches:
        # Parse options
        options_pattern = r'[Aa]\)\s*(.*?)\s*[Bb]\)\s*(.*?)\s*[Cc]\)\s*(.*?)\s*[Dd]\)\s*(.*?)(?:\n|$)'
        options_match = re.search(options_pattern, q_text, re.DOTALL)
        
        if options_match:
            question_text = q_text[:options_match.start()].strip()
            options = {
                "A": options_match.group(1).strip(),
                "B": options_match.group(2).strip(),
                "C": options_match.group(3).strip(),
                "D": options_match.group(4).strip()
            }
            
            questions.append({
                "number": int(number),
                "question": question_text,
                "options": options
            })
    
    # Parse answer key
    if answers_text:
        answer_pattern = r'(\d+)\s*[-–]?\s*([A-Da-d])'
        answer_matches = re.findall(answer_pattern, answers_text)
        
        for num, ans in answer_matches:
            answer_key[num] = ans.upper()
    
    return questions, answer_key

# ============================================================================
# HEALTH ENDPOINTS
# ============================================================================
@app.get("/health")
async def health():
    return {"status": "healthy", "timestamp": datetime.utcnow().isoformat()}

@app.get("/ping")
async def ping():
    return {"message": "pong"}

@app.get("/test-user-query")
async def test_user_query(db: Session = Depends(get_db)):
    """Test endpoint to check if user query works"""
    try:
        user = db.query(User).first()
        if user:
            return {
                "status": "success",
                "user": {
                    "id": user.id,
                    "email": user.email,
                    "name": getattr(user, 'name', None),
                    "surname": getattr(user, 'surname', None),
                    "tier": getattr(user, 'tier', 'free'),
                    "is_admin": bool(getattr(user, 'is_admin', 0)) if getattr(user, 'is_admin', 0) else False
                }
            }
        else:
            return {"status": "no users found"}
    except Exception as e:
        import traceback
        return {
            "status": "error",
            "error": str(e),
            "traceback": traceback.format_exc()
        }

@app.options("/{path:path}")
async def options_handler(path: str):
    """Handle OPTIONS requests for CORS preflight"""
    return JSONResponse(content={"status": "ok"})

# ============================================================================
# AUTH ENDPOINTS
# ============================================================================
@app.post("/auth/register")
async def register(user_data: UserRegister, db: Session = Depends(get_db)):
    # Check if user exists
    existing = db.query(User).filter(User.email == user_data.email).first()
    if existing:
        raise HTTPException(status_code=400, detail="Email already registered")
    
    # Create user
    user = User(
        email=user_data.email,
        password_hash=hash_password(user_data.password),
        name=user_data.name,
        surname=user_data.surname,
        tier="free"
    )
    db.add(user)
    db.commit()
    db.refresh(user)
    
    return {"id": user.id, "email": user.email, "name": user.name, "surname": user.surname, "tier": user.tier}

@app.post("/auth/login")
async def login(credentials: UserLogin, response: Response, db: Session = Depends(get_db)):
    try:
        print(f"[AUTH] Login attempt for email: {credentials.email}")
        
        # Query user - use try/except to catch any database errors
        try:
            user = db.query(User).filter(User.email == credentials.email).first()
        except Exception as db_error:
            print(f"[AUTH ERROR] Database query failed: {db_error}")
            import traceback
            traceback.print_exc()
            raise HTTPException(status_code=500, detail="Database error during login")
        
        if not user:
            print(f"[AUTH] User not found: {credentials.email}")
            raise HTTPException(status_code=401, detail="Invalid credentials")
        
        print(f"[AUTH] User found: {user.id}, email: {user.email}")
        
        # Verify password
        try:
            if not user.password_hash:
                print(f"[AUTH] User {user.id} has no password hash")
                raise HTTPException(status_code=401, detail="Invalid credentials")
            
            if not verify_password(credentials.password, user.password_hash):
                print(f"[AUTH] Password verification failed for user {user.id}")
                raise HTTPException(status_code=401, detail="Invalid credentials")
        except HTTPException:
            raise
        except Exception as pwd_error:
            print(f"[AUTH ERROR] Password verification error: {pwd_error}")
            raise HTTPException(status_code=401, detail="Invalid credentials")
        
        print(f"[AUTH] Password verified for user {user.id}")
        
        # Create tokens
        try:
            access_token = create_token(user.id, timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES))
            refresh_token = create_token(user.id, timedelta(days=REFRESH_TOKEN_EXPIRE_DAYS))
            print(f"[AUTH] Tokens created for user {user.id}")
        except Exception as token_error:
            print(f"[AUTH ERROR] Token creation failed: {token_error}")
            import traceback
            traceback.print_exc()
            raise HTTPException(status_code=500, detail="Failed to create authentication tokens")
        
        # Set HTTP-only cookies for security
        try:
            response.set_cookie(
                key="access_token",
                value=access_token,
                path="/",
                httponly=True,
                secure=COOKIE_SECURE,
                samesite=COOKIE_SAMESITE,
                max_age=ACCESS_TOKEN_EXPIRE_MINUTES * 60
            )
            response.set_cookie(
                key="refresh_token",
                value=refresh_token,
                path="/",
                httponly=True,
                secure=COOKIE_SECURE,
                samesite=COOKIE_SAMESITE,
                max_age=REFRESH_TOKEN_EXPIRE_DAYS * 24 * 60 * 60
            )
            print(f"[AUTH] Cookies set: secure={COOKIE_SECURE}, samesite={COOKIE_SAMESITE}")
        except Exception as cookie_error:
            print(f"[AUTH WARNING] Cookie setting failed: {cookie_error}")
            # Continue anyway - tokens are still returned in response body
        
        # Build user response data - use direct attribute access with fallbacks
        try:
            # Get basic user info
            user_id = user.id
            user_email = user.email
            
            # Get optional fields with safe defaults
            user_tier = 'free'
            try:
                if hasattr(user, 'tier') and user.tier:
                    user_tier = user.tier
            except:
                pass
            
            user_name = None
            try:
                if hasattr(user, 'name'):
                    user_name = user.name
            except:
                pass
            
            user_surname = None
            try:
                if hasattr(user, 'surname'):
                    user_surname = user.surname
            except:
                pass
            
            # Get is_admin safely
            is_admin = False
            try:
                if hasattr(user, 'is_admin'):
                    is_admin_val = user.is_admin
                    if is_admin_val is not None:
                        is_admin = bool(int(is_admin_val))
            except:
                pass
            
            user_data = {
                "id": user_id,
                "email": user_email,
                "tier": user_tier,
                "name": user_name,
                "surname": user_surname,
                "is_admin": is_admin
            }
            
            print(f"[AUTH] User data built: {user_data}")
            
        except Exception as user_data_error:
            print(f"[AUTH ERROR] Error building user data: {user_data_error}")
            import traceback
            traceback.print_exc()
            # Fallback to minimal user data
            user_data = {
                "id": user.id,
                "email": user.email,
                "tier": 'free',
                "is_admin": False
            }
        
        response_data = {
            "access_token": access_token,
            "refresh_token": refresh_token,
            "token_type": "bearer",
            "user": user_data
        }
        
        print(f"[AUTH] Login successful for user {user.id}")
        return response_data
        
    except HTTPException:
        # Re-raise HTTP exceptions (like 401)
        raise
    except Exception as e:
        # Log the error for debugging
        import traceback
        error_details = str(e)
        traceback_str = traceback.format_exc()
        print(f"[AUTH ERROR] Login failed with exception: {error_details}")
        print(f"[AUTH ERROR] Full traceback:\n{traceback_str}")
        raise HTTPException(status_code=500, detail="Login failed. Please try again.")

@app.post("/auth/refresh")
async def refresh(request: Request, response: Response, refresh_data: dict = None):
    # Try to get refresh token from cookie first, then from request body
    refresh_token = request.cookies.get("refresh_token")
    
    if not refresh_token and refresh_data:
        refresh_token = refresh_data.get("refresh_token")
    
    if not refresh_token:
        raise HTTPException(status_code=401, detail="Missing refresh token")
    
    try:
        payload = jwt.decode(refresh_token, SECRET_KEY, algorithms=[ALGORITHM])
        user_id_str = payload.get("sub")
        user_id = int(user_id_str)  # Convert string to int
        
        access_token = create_token(user_id, timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES))
        new_refresh = create_token(user_id, timedelta(days=REFRESH_TOKEN_EXPIRE_DAYS))
        
        # Set new cookies
        response.set_cookie(
            key="access_token",
            value=access_token,
            path="/",
            httponly=True,
            secure=COOKIE_SECURE,
            samesite=COOKIE_SAMESITE,
            max_age=ACCESS_TOKEN_EXPIRE_MINUTES * 60
        )
        response.set_cookie(
            key="refresh_token",
            value=new_refresh,
            path="/",
            httponly=True,
            secure=COOKIE_SECURE,
            samesite=COOKIE_SAMESITE,
            max_age=REFRESH_TOKEN_EXPIRE_DAYS * 24 * 60 * 60
        )
        
        return {
            "access_token": access_token,
            "refresh_token": new_refresh,
            "token_type": "bearer"
        }
    except JWTError:
        raise HTTPException(status_code=401, detail="Invalid refresh token")

@app.post("/auth/logout")
async def logout(response: Response):
    # Clear cookies
    response.delete_cookie(key="access_token", path="/", samesite=COOKIE_SAMESITE, secure=COOKIE_SECURE)
    response.delete_cookie(key="refresh_token", path="/", samesite=COOKIE_SAMESITE, secure=COOKIE_SECURE)
    return {"message": "Logged out successfully"}

@app.get("/me")
async def get_me(current_user: User = Depends(get_current_user), db: Session = Depends(get_db)):
    try:
        # Get today's usage
        today = date.today()
        usage_data = {}
        
        for kind in ["exam", "explain", "chat", "upload"]:
            try:
                usage = db.query(Usage).filter(
                    Usage.user_id == current_user.id,
                    Usage.kind == kind,
                    Usage.date == today
                ).first()
                usage_data[kind] = usage.count if usage else 0
            except Exception as e:
                # If usage table has issues, set to 0 and continue
                print(f"[WARNING] Failed to get usage for {kind}: {e}")
                usage_data[kind] = 0
        
        # Safely get is_admin - handle case where column might not exist or be None
        # is_admin is stored as integer (0 or 1)
        is_admin_value = getattr(current_user, 'is_admin', 0) or 0
        is_admin = bool(is_admin_value) if is_admin_value else False
        
        return {
            "id": current_user.id,
            "email": current_user.email,
            "name": getattr(current_user, 'name', None) if hasattr(current_user, 'name') else None,
            "surname": getattr(current_user, 'surname', None) if hasattr(current_user, 'surname') else None,
            "tier": getattr(current_user, 'tier', 'free') if hasattr(current_user, 'tier') else 'free',
            "is_admin": is_admin,
            "usage": usage_data
        }
    except Exception as e:
        # Log the error and return a 500 with details
        import traceback
        error_details = str(e)
        traceback_str = traceback.format_exc()
        print(f"[ERROR] Error in /me endpoint: {error_details}\n{traceback_str}")
        raise HTTPException(status_code=500, detail=f"Failed to get user info: {error_details}")

# ============================================================================
# FILE UPLOAD ENDPOINT
# ============================================================================
@app.post("/upload")
async def upload_files(
    request: Request,
    files: List[UploadFile] = File(...),
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_db)
):
    # For demo purposes, allow uploads without authentication but with limits
    if not current_user:
        # Anonymous uploads limited to 1 file
        if len(files) > 1:
            raise HTTPException(status_code=403, detail="Please login to upload multiple files")
    else:
        # Check quota for authenticated users
        if not check_quota(db, current_user, "upload"):
            raise HTTPException(status_code=403, detail="Upload quota exceeded. Please upgrade to Premium.")
    
    results = []
    
    for file in files:
        file_bytes = await file.read()
        
        # Extract text content from file
        try:
            text_content = extract_text_from_file(
                file_bytes, 
                file.filename,
                file.content_type or "application/octet-stream"
            )
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to process {file.filename}: {str(e)}")
        
        # Generate a simple file_id for tracking
        import hashlib
        file_id = f"file-{hashlib.md5(file_bytes).hexdigest()[:16]}"
        
        # Save to database if user is authenticated, otherwise use in-memory storage
        if current_user:
            upload = Upload(
                user_id=current_user.id,
                file_id=file_id,
                filename=file.filename,
                mime=file.content_type or "application/octet-stream",
                size=len(file_bytes),
                content=text_content
            )
            db.add(upload)
        else:
            # Store in memory for anonymous users
            file_content_store[file_id] = {
                "filename": file.filename,
                "content": text_content,
                "mime": file.content_type or "application/octet-stream"
            }
        
        results.append({
            "file_id": file_id,
            "filename": file.filename,
            "mime": file.content_type,
            "size": len(file_bytes),
            "content_length": len(text_content)
        })
    
    if current_user:
        db.commit()
        increment_usage(db, current_user, "upload")
    
    return results

# ============================================================================
# GROUNDED ENDPOINTS
# ============================================================================
@app.post("/summarize-from-files")
async def summarize_from_files(
    request: Request,
    req: SummarizeRequest,
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_db)
):
    """
    Advanced summarization endpoint with:
    - Plan-based limits (free/standard/premium)
    - Map-reduce for large documents
    - SHA256 caching for deduplication
    - Token estimation and adaptive output sizing
    """
    rate_limit(request)
    
    # Import new modular services
    from app.config import PLAN_LIMITS, ALLOWED_EXTS, TOKEN_PER_CHAR
    from app.utils.files import (
        ext_ok, pdf_page_count, approx_tokens_from_text_len,
        sha256_bytes, choose_max_output_tokens, validate_mime_type, basic_antivirus_check
    )
    from app.services.cache import get_cached, set_cached
    from app.services.summary import map_reduce_summary, summarize_no_files
    import json
    import hashlib
    
    # Determine user's plan
    plan = getattr(current_user, "tier", "free") if current_user else "free"
    # Map "premium" to "pro" if needed
    if plan == "premium":
        plan = "pro"
    limits = PLAN_LIMITS.get(plan, PLAN_LIMITS["standard"])
    
    # Require either file_ids or prompt
    if not req.file_ids and not req.prompt:
        raise HTTPException(status_code=400, detail="Please provide either files or a prompt.")
    
    # ========== CASE 1: No files, just prompt ==========
    if not req.file_ids:
        try:
            language = req.language or "en"
            result_json = summarize_no_files(
                topic=req.prompt,
                language=language,
                out_cap=limits.max_output_cap
            )
            
            # Parse and return
            result_json = result_json.strip()
            if result_json.startswith('```'):
                lines = result_json.split('\n')
                if lines[0].strip().startswith('```'):
                    lines = lines[1:]
                if lines and lines[-1].strip() == '```':
                    lines = lines[:-1]
                result_json = '\n'.join(lines)
            
            try:
                result = json.loads(result_json)
            except json.JSONDecodeError:
                import re
                json_match = re.search(r'\{.*\}', result_json, re.DOTALL)
                if json_match:
                    try:
                        result = json.loads(json_match.group(0))
                    except:
                        result = {
                            "summary": {
                                "title": "Summary",
                                "sections": [{"heading": "Content", "bullets": [result_json]}]
                            },
                            "citations": []
                        }
                else:
                    result = {
                        "summary": {
                            "title": "Summary",
                            "sections": [{"heading": "Content", "bullets": [result_json]}]
                        },
                        "citations": []
                    }
            
            # Ensure result has correct structure: {summary: {...}, citations: [...]}
            if "summary" not in result:
                # Check for core_concepts (common in no-files prompts)
                if "core_concepts" in result and not result.get("sections"):
                    print("[SUMMARY NO-FILES] Found core_concepts, transforming to sections...")
                    sections = []
                    for idx, concept in enumerate(result.get("core_concepts", [])):
                        if isinstance(concept, dict):
                            concept_term = concept.get("concept", concept.get("term", f"Concept {idx + 1}"))
                            concept_explanation = concept.get("explanation", "")
                            concept_examples = concept.get("examples", [])
                            
                            definition = concept.get("definition", "")
                            if not definition and concept_explanation:
                                sentences = concept_explanation.split(".")
                                definition = sentences[0].strip() + "." if sentences[0].strip() else concept_explanation[:100]
                            
                            example = ""
                            if concept_examples:
                                if isinstance(concept_examples[0], dict):
                                    example = concept_examples[0].get("example", "")
                                elif isinstance(concept_examples[0], str):
                                    example = concept_examples[0]
                            
                            sections.append({
                                "heading": concept_term,
                                "concepts": [{
                                    "term": concept_term,
                                    "definition": definition or concept_term,
                                    "explanation": concept_explanation or definition,
                                    "example": example,
                                    "key_points": concept.get("key_points", []) if isinstance(concept.get("key_points"), list) else []
                                }]
                            })
                    
                    result = {
                        "summary": {
                            "title": result.get("title", "Summary"),
                            "overview": result.get("overview", ""),
                            "learning_objectives": result.get("learning_objectives", []),
                            "sections": sections,
                            "formula_sheet": result.get("formula_sheet", []),
                            "glossary": result.get("glossary", [])
                        },
                        "citations": []
                    }
                elif "title" in result or "sections" in result:
                    # It's the summary object itself, wrap it
                    result = {
                        "summary": result,
                        "citations": result.get("citations", [])
                    }
                else:
                    # Unknown structure, create default - DO NOT put the entire result as a string!
                    print(f"[SUMMARY NO-FILES] Unknown structure: {list(result.keys())}")
                    result = {
                        "summary": {
                            "title": "Summary",
                            "overview": "Content generated successfully",
                            "sections": [{"heading": "Content", "bullets": ["Please check the console for detailed structure information."]}]
                        },
                        "citations": []
                    }
            elif "citations" not in result:
                result["citations"] = []
            
            return result
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Summary generation failed: {str(e)}")
    
    # ========== CASE 2: Files provided ==========
    
    # Fetch and validate files
    files_data = []  # [(filename, content_bytes, text_content)]
    total_mb = 0.0
    total_pages = 0
    pdf_count = 0
    
    for file_id in req.file_ids:
        # Try database first
        upload = db.query(Upload).filter(Upload.file_id == file_id).first()
        if not upload:
            # Try in-memory storage for anonymous users
            if file_id in file_content_store:
                file_info = file_content_store[file_id]
                filename = file_info["filename"]
                text_content = file_info["content"]
                # Reconstruct bytes for validation (approximate)
                content_bytes = text_content.encode("utf-8", errors="ignore")
            else:
                raise HTTPException(status_code=404, detail=f"File not found: {file_id}")
        else:
            filename = upload.filename
            text_content = upload.content  # Already extracted text
            content_bytes = text_content.encode("utf-8", errors="ignore")
        
        # File type validation
        if not ext_ok(filename, ALLOWED_EXTS):
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {filename}")
        
        # Size tracking
        size_mb = len(content_bytes) / (1024 * 1024)
        total_mb += size_mb
        
        # PDF page count
        if filename.lower().endswith(".pdf"):
            pdf_count += 1
            # For PDF, we'd need original bytes, but we have text already
            # Skip page count validation if we only have text
            # In production, store original bytes too
        
        # Basic security checks
        if not basic_antivirus_check(content_bytes):
            raise HTTPException(status_code=400, detail=f"File failed security check: {filename}")
        
        files_data.append((filename, content_bytes, text_content))
    
    # ========== LIMIT CHECKS ==========
    if len(files_data) > limits.max_files_total:
        raise HTTPException(
            status_code=400,
            detail=f"Too many files. Your plan allows {limits.max_files_total} files. Upgrade for more."
        )
    
    if pdf_count > limits.max_pdfs:
        raise HTTPException(
            status_code=400,
            detail=f"Too many PDFs. Your plan allows {limits.max_pdfs} PDFs. Upgrade for more."
        )
    
    if total_mb > limits.max_total_mb:
        raise HTTPException(
            status_code=400,
            detail=f"Total file size ({total_mb:.1f} MB) exceeds your plan limit ({limits.max_total_mb} MB)."
        )
    
    # Merge all text content
    merged_text = "\n\n".join([text for _, _, text in files_data if text])
    
    # Token estimation
    estimated_tokens = approx_tokens_from_text_len(len(merged_text), TOKEN_PER_CHAR)
    
    # Hard cap to prevent abuse (5x plan limit)
    HARD_CAP = limits.max_input_tokens * 5
    if estimated_tokens > HARD_CAP:
        raise HTTPException(
            status_code=400,
            detail=f"Content extremely large (~{estimated_tokens} tokens). Maximum is {HARD_CAP} tokens. Please split the upload."
        )
    
    # If content is large (>50% of plan limit), force map-reduce
    force_map_reduce = estimated_tokens > (limits.max_input_tokens // 2)
    
    # Use plan's max output cap directly (always use max for quality)
    # Adaptive sizing was causing too-low limits
    out_cap = limits.max_output_cap
    
    print(f"[SUMMARY CONFIG] Plan: {plan}, max_output_cap: {limits.max_output_cap}, using: {out_cap}")
    
    # ========== CACHE CHECK ==========
    # Create cache key from: plan, language, prompt, file hashes, out_cap, model
    from app.config import OPENAI_MODEL
    cache_key_data = {
        "plan": plan,
        "language": req.language or "en",
        "prompt": (req.prompt or "")[:1000],  # Limit prompt length in cache key
        "out_cap": out_cap,
        "model": OPENAI_MODEL,
        "file_hashes": [sha256_bytes(content) for _, content, _ in files_data]
    }
    cache_key = hashlib.sha256(
        json.dumps(cache_key_data, sort_keys=True).encode()
    ).hexdigest()
    
    # Check cache
    cached_result = get_cached(cache_key, db)
    if cached_result:
        print(f"[CACHE HIT] Returning cached summary for {cache_key[:12]}...")
        return json.loads(cached_result)
    
    # ========== GENERATE SUMMARY ==========
    print(f"[CACHE MISS] Generating new summary (estimated={estimated_tokens} tokens, out_cap={out_cap}, force_map_reduce={force_map_reduce})...")
    
    try:
        from app.utils.json_helpers import parse_json_robust, create_error_response
        from app.utils.quality import enforce_exam_ready, validate_summary_completeness, calculate_comprehensive_quality_score
        from app.services.summary import quality_score_legacy
        
        language = req.language or "en"
        additional_instructions = req.prompt or ""
        
        # Track generation start time
        generation_start = time.time()
        
        # Use map-reduce pipeline
        result_json = map_reduce_summary(
            full_text=merged_text,
            language=language,
            additional_instructions=additional_instructions,
            out_cap=out_cap,
            force_chunking=force_map_reduce
        )
        
        # Parse JSON with robust error handling
        try:
            result = parse_json_robust(result_json)
            print("[SUMMARY] JSON parsed successfully")
            
            # Ensure result has the correct structure: {summary: {...}, citations: [...]}
            # If result doesn't have 'summary' key, it might be the summary object itself
            if "summary" not in result:
                print("[SUMMARY] Result doesn't have 'summary' key, checking structure...")
                print(f"[SUMMARY] Result keys: {list(result.keys())}")
                
                # Check if it has learning_objectives or core_concepts (common AI output formats)
                has_learning_objectives = "learning_objectives" in result
                has_core_concepts = "core_concepts" in result
                has_title = "title" in result
                has_sections = "sections" in result
                
                # PRIORITY 1: Transform core_concepts into sections if needed (most common issue)
                if has_core_concepts and not has_sections:
                    print("[SUMMARY] Found core_concepts, transforming to sections structure...")
                    # Convert core_concepts to sections format
                    sections = []
                    for idx, concept in enumerate(result.get("core_concepts", [])):
                        # If concept is a dict with 'concept' key (old format), transform it
                        if isinstance(concept, dict):
                            concept_term = concept.get("concept", concept.get("term", f"Concept {idx + 1}"))
                            concept_explanation = concept.get("explanation", "")
                            concept_examples = concept.get("examples", [])
                            
                            # Extract definition from explanation if not provided
                            definition = concept.get("definition", "")
                            if not definition and concept_explanation:
                                # Use first sentence as definition
                                sentences = concept_explanation.split(".")
                                definition = sentences[0].strip() + "." if sentences[0].strip() else concept_explanation[:100]
                            
                            # Extract example
                            example = ""
                            if concept_examples:
                                if isinstance(concept_examples[0], dict):
                                    example = concept_examples[0].get("example", "")
                                elif isinstance(concept_examples[0], str):
                                    example = concept_examples[0]
                            
                            # Create a section for this concept
                            section = {
                                "heading": concept_term,
                                "concepts": [{
                                    "term": concept_term,
                                    "definition": definition or concept_term,
                                    "explanation": concept_explanation or definition,
                                    "example": example,
                                    "key_points": concept.get("key_points", []) if isinstance(concept.get("key_points"), list) else []
                                }]
                            }
                            sections.append(section)
                        else:
                            # If concept is a string or other format, create a basic section
                            sections.append({
                                "heading": f"Concept {idx + 1}",
                                "concepts": [{
                                    "term": str(concept)[:50],
                                    "definition": str(concept),
                                    "explanation": str(concept),
                                    "key_points": []
                                }]
                            })
                    
                    # Create proper summary structure
                    result = {
                        "summary": {
                            "title": result.get("title", "Summary"),
                            "overview": result.get("overview", ""),
                            "learning_objectives": result.get("learning_objectives", []),
                            "sections": sections,
                            "formula_sheet": result.get("formula_sheet", []),
                            "glossary": result.get("glossary", [])
                        },
                        "citations": result.get("citations", [])
                    }
                    print(f"[SUMMARY] Transformed to sections structure with {len(sections)} sections")
                
                # PRIORITY 2: Check if it looks like a summary object (has title or sections)
                elif has_title or has_sections:
                    # It's the summary object itself, wrap it
                    print("[SUMMARY] Result is summary object, wrapping it")
                    result = {
                        "summary": result,
                        "citations": result.get("citations", [])
                    }
                else:
                    # Unknown structure - check if it's a dict that might contain the summary
                    print(f"[SUMMARY] Unknown result structure: {list(result.keys())}")
                    # Try to see if it's a string representation of the summary
                    if len(result) == 1 and isinstance(list(result.values())[0], str):
                        # Might be a stringified JSON
                        try:
                            import json
                            parsed = json.loads(list(result.values())[0])
                            if "summary" in parsed or "sections" in parsed or "learning_objectives" in parsed:
                                result = parsed
                                if "summary" not in result:
                                    result = {"summary": result, "citations": []}
                        except:
                            pass
                    
                    # If still not valid, create error response
                    if "summary" not in result:
                        print(f"[SUMMARY] Creating error response for invalid structure")
                        result = create_error_response(
                            f"Unexpected response structure. Keys: {list(result.keys())}. Expected: summary with sections, or learning_objectives/core_concepts.",
                            len(result_json)
                        )
            else:
                # Ensure citations exist
                if "citations" not in result:
                    result["citations"] = []
                    
        except ValueError as e:
            print(f"[SUMMARY] All JSON parse attempts failed: {e}")
            result = create_error_response(
                "Failed to parse AI response. This may be due to response format issues.",
                len(result_json)
            )
        
        # Calculate comprehensive quality score - ensure summary exists
        if "summary" not in result:
            print("[SUMMARY ERROR] Result still doesn't have 'summary' key after normalization")
            result = create_error_response("Internal error: summary structure invalid", 0)
        
        quality_metrics = calculate_comprehensive_quality_score(result)
        score = quality_metrics.get("final_ready_score", 0.5)
        is_final_ready = quality_metrics.get("is_final_ready", False)
        
        print(f"[QUALITY METRICS] Final-ready score: {score}/1.0 (target: 0.90+)")
        print(f"[QUALITY METRICS] Coverage: {quality_metrics.get('coverage_score', 0)}, " +
              f"Numeric density: {quality_metrics.get('numeric_density', 0)}, " +
              f"Formula completeness: {quality_metrics.get('formula_completeness', 0)}, " +
              f"Citation depth: {quality_metrics.get('citation_depth', 0)}, " +
              f"Readability: {quality_metrics.get('readability_score', 0)}")
        print(f"[QUALITY METRICS] Domain: {quality_metrics.get('domain', 'unknown')}, " +
              f"Is final-ready: {is_final_ready}")
        
        # Enforce exam-ready quality standards
        result = enforce_exam_ready(result, detected_themes=None)
        
        # POST-PROCESSING VALIDATION
        from app.utils.quality import create_self_repair_prompt, validate_summary_completeness, validate_and_enhance_quality
        from app.services.summary import call_openai, SYSTEM_PROMPT
        
        # Step 1: Enhance and validate
        result, repair_prompts = validate_and_enhance_quality(result)
        print(f"[POST-PROCESSING] Auto-enhancements applied, {len(repair_prompts)} repair prompts generated")
        
        # Step 2: Completeness check
        warnings, needs_repair = validate_summary_completeness(result)
        
        if warnings:
            print(f"[SUMMARY QUALITY] Warnings ({len(warnings)}): {warnings}")
        
        # Track self-repair metrics
        self_repair_triggered = False
        self_repair_improvement = None
        
        # Step 3: Trigger self-repair if needed
        if (repair_prompts or needs_repair) and score < 0.7:
            self_repair_triggered = True
            
            # Combine all repair prompts
            combined_repairs = []
            if repair_prompts:
                combined_repairs.extend(repair_prompts)
            if needs_repair and warnings:
                combined_repairs.append(create_self_repair_prompt(result, warnings, language))
            
            repair_instruction = "\n\n---\n\n".join(combined_repairs)
            print(f"[SELF-REPAIR] Triggering repair (score: {score}, {len(combined_repairs)} issues)")
            
            try:
                repaired_json = call_openai(
                    system_prompt=SYSTEM_PROMPT,
                    user_prompt=f"Fix the following issues in this summary:\n\n{repair_instruction}",
                    max_output_tokens=min(out_cap, 8000),
                    retry_on_length=False
                )
                
                # Try to parse repaired output
                try:
                    repaired_result = parse_json_robust(repaired_json)
                    
                    # Ensure repaired_result has correct structure
                    if "summary" not in repaired_result:
                        if "title" in repaired_result or "sections" in repaired_result:
                            repaired_result = {
                                "summary": repaired_result,
                                "citations": repaired_result.get("citations", [])
                            }
                        else:
                            print(f"[SELF-REPAIR] Repaired result has invalid structure, keeping original")
                            repaired_result = result
                    elif "citations" not in repaired_result:
                        repaired_result["citations"] = []
                    
                    # Check if repair improved quality
                    repaired_metrics = calculate_comprehensive_quality_score(repaired_result)
                    repaired_score = repaired_metrics.get("final_ready_score", 0.5)
                    print(f"[SELF-REPAIR] Score after repair: {repaired_score}/1.0")
                    
                    if repaired_score > score:
                        self_repair_improvement = repaired_score - score
                        print(f"[SELF-REPAIR] Accepted (improvement: +{self_repair_improvement:.2f})")
                        result = repaired_result
                        score = repaired_score
                    else:
                        print(f"[SELF-REPAIR] Rejected (no improvement)")
                except Exception as e:
                    print(f"[SELF-REPAIR] Parse failed: {e}, keeping original")
            except Exception as e:
                print(f"[SELF-REPAIR] Failed: {e}, keeping original")
        
        # Calculate generation time and metrics
        generation_time = time.time() - generation_start
        
        # Extract quality metrics for telemetry
        summary = result.get("summary", {})
        num_concepts = sum(len(sec.get("concepts", [])) for sec in summary.get("sections", []))
        num_formulas = len(summary.get("formula_sheet", []))
        num_exam_questions = 0
        num_glossary = len(summary.get("glossary", []))
        
        # Record telemetry (non-blocking)
        try:
            from app.services.telemetry import record_summary_quality
            from app.services.summary import detect_domain
            
            # Detect domain from original text
            domain = detect_domain(merged_text)
            
            # Estimate total tokens used
            total_tokens_used = estimated_tokens + out_cap
            
            record_summary_quality(
                db=db,
                request_hash=cache_key,
                user_id=current_user.id if current_user else None,
                plan=plan,
                domain=domain,
                language=language,
                input_tokens=estimated_tokens,
                num_chunks=len(files_data),
                quality_score=score,
                num_concepts=num_concepts,
                num_formulas=num_formulas,
                num_exam_questions=num_exam_questions,
                num_glossary_terms=num_glossary,
                self_repair_triggered=self_repair_triggered,
                self_repair_improvement=self_repair_improvement,
                total_tokens_used=total_tokens_used,
                generation_time_seconds=generation_time,
                warnings=warnings,
                coverage_score=quality_metrics.get('coverage_score'),
                numeric_density=quality_metrics.get('numeric_density'),
                formula_completeness=quality_metrics.get('formula_completeness'),
                citation_depth=quality_metrics.get('citation_depth'),
                readability_score=quality_metrics.get('readability_score'),
                is_final_ready=is_final_ready
            )
        except Exception as telemetry_error:
            print(f"[TELEMETRY WARNING] Failed to record: {telemetry_error}")
        
        # Final validation: ensure result has correct structure before returning
        if "summary" not in result:
            print("[SUMMARY ERROR] Final validation failed: result missing 'summary' key")
            result = create_error_response("Internal error: invalid summary structure", 0)
        elif "citations" not in result:
            result["citations"] = []
        
        # Cache the result (only if not error)
        if "error" not in result.get("summary", {}).get("title", "").lower():
            set_cached(cache_key, json.dumps(result), db)
        
        print(f"[SUMMARY] Returning result with structure: summary={bool(result.get('summary'))}, citations={bool(result.get('citations'))}")
        return result
        
    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"[SUMMARY ERROR] {error_trace}")
        raise HTTPException(status_code=500, detail=f"Summary generation failed: {str(e)}")

@app.post("/flashcards-from-files")
async def flashcards_from_files(
    request: Request,
    req: FlashcardsRequest,
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_db)
):
    rate_limit(request)
    
    # Require either file_ids or prompt
    if not req.file_ids and not req.prompt:
        raise HTTPException(status_code=400, detail="Please provide either files or a prompt.")
    
    # Fetch file contents from database or in-memory storage
    file_contents = []
    if req.file_ids:
        for file_id in req.file_ids:
            # Try database first
            upload = db.query(Upload).filter(Upload.file_id == file_id).first()
            if upload and upload.content:
                file_contents.append(upload.content)
            # Try in-memory storage for anonymous users
            elif file_id in file_content_store:
                file_contents.append(file_content_store[file_id]["content"])
    
    additional_instructions = ""
    if req.prompt:
        additional_instructions = f"\n\nADDITIONAL USER INSTRUCTIONS:\n{req.prompt}\n\nMake sure to follow these instructions while creating flashcards."
    
    language_instruction = ""
    if req.language == "tr":
        language_instruction = "\n\nIMPORTANT: Generate ALL flashcards in TURKISH language."
    elif req.language == "en":
        language_instruction = "\n\nIMPORTANT: Generate ALL flashcards in ENGLISH language."
    
    system_prompt = """You are a study assistant. Create helpful flashcards from the document content or user topics."""
    
    if file_contents:
        user_prompt = f"""Create {req.count} flashcards from the documents. Extract key concepts, questions, or important information.
{language_instruction}
{additional_instructions}

Output as JSON:
{{
  "deck": "{req.deck_name}",
  "cards": [
    {{"type": "{req.style}", "front": "Question or term", "back": "Answer or definition", "source": {{"file_id": "doc", "evidence": "relevant text"}}}}
  ]
}}"""
    else:
        # No files, just prompt
        user_prompt = f"""Create {req.count} flashcards about the topic requested by the user.
{language_instruction}

USER REQUEST:
{req.prompt}

Output as JSON:
{{
  "deck": "{req.deck_name}",
  "cards": [
    {{"type": "{req.style}", "front": "Question or term", "back": "Answer or definition"}}
  ]
}}"""
    
    try:
        user_id = current_user.id if current_user else None
        response_text = call_openai_with_context(
            file_contents, 
            f"{system_prompt}\n\n{user_prompt}", 
            temperature=0.0,
            user_id=user_id,
            endpoint="flashcards",
            db=db
        )
        
        import json
        
        # Clean the response to extract JSON
        response_text = response_text.strip()
        
        # Try to find JSON in response
        if response_text.startswith('```'):
            # Remove code blocks
            lines = response_text.split('\n')
            response_text = '\n'.join([l for l in lines if not l.strip().startswith('```')])
        
        try:
            result = json.loads(response_text)
        except:
            # If still fails, try to extract JSON from text
            import re
            json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
            if json_match:
                result = json.loads(json_match.group(0))
            else:
                # Fallback: create cards from the response
                result = {
                    "deck": req.deck_name,
                    "cards": [{"type": req.style, "front": "Content", "back": response_text, "source": {"file_id": "", "evidence": ""}}]
                }
        
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/truefalse-from-files")
async def truefalse_from_files(
    request: Request,
    req: TrueFalseRequest,
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_db)
):
    rate_limit(request)
    
    # Require either file_ids or prompt
    if not req.file_ids and not req.prompt:
        raise HTTPException(status_code=400, detail="Please provide either files or a prompt.")
    
    # Fetch file contents from database or in-memory storage
    file_contents = []
    if req.file_ids:
        for file_id in req.file_ids:
            # Try database first
            upload = db.query(Upload).filter(Upload.file_id == file_id).first()
            if upload and upload.content:
                file_contents.append(upload.content)
            # Try in-memory storage for anonymous users
            elif file_id in file_content_store:
                file_contents.append(file_content_store[file_id]["content"])
    
    additional_instructions = ""
    if req.prompt:
        additional_instructions = f"\n\nADDITIONAL USER INSTRUCTIONS:\n{req.prompt}\n\nMake sure to follow these instructions while creating True/False statements."
    
    language_instruction = ""
    if req.language == "tr":
        language_instruction = "\n\nIMPORTANT: Generate ALL True/False statements in TURKISH language."
    elif req.language == "en":
        language_instruction = "\n\nIMPORTANT: Generate ALL True/False statements in ENGLISH language."
    
    system_prompt = """You are a study assistant. Create True/False statement cards from document content or user topics."""
    
    if file_contents:
        user_prompt = f"""Create {req.count} True/False statement cards from the documents. Each statement should be a factual claim that can be verified as true or false based on the document content.
{language_instruction}
{additional_instructions}

IMPORTANT RULES:
- Each statement should be a clear, factual claim
- Statements should be based ONLY on the document content
- Mix true and false statements (approximately 50/50 or as appropriate)
- For false statements, make them plausible but incorrect based on the content
- Each statement should be standalone and verifiable

Output as JSON:
{{
  "cards": [
    {{"statement": "Statement text here", "answer": true, "explanation": "Brief explanation why this is true/false"}},
    {{"statement": "Another statement", "answer": false, "explanation": "Brief explanation why this is false"}}
  ]
}}

Return ONLY valid JSON, no markdown code blocks, no extra text."""
    else:
        # No files, just prompt
        user_prompt = f"""Create {req.count} True/False statement cards about the topic requested by the user.
{language_instruction}

USER REQUEST:
{req.prompt}

IMPORTANT RULES:
- Each statement should be a clear, factual claim
- Mix true and false statements (approximately 50/50 or as appropriate)
- For false statements, make them plausible but incorrect
- Each statement should be standalone and verifiable

Output as JSON:
{{
  "cards": [
    {{"statement": "Statement text here", "answer": true, "explanation": "Brief explanation why this is true/false"}},
    {{"statement": "Another statement", "answer": false, "explanation": "Brief explanation why this is false"}}
  ]
}}

Return ONLY valid JSON, no markdown code blocks, no extra text."""
    
    try:
        user_id = current_user.id if current_user else None
        response_text = call_openai_with_context(
            file_contents, 
            f"{system_prompt}\n\n{user_prompt}", 
            temperature=0.0,
            user_id=user_id,
            endpoint="truefalse",
            db=db
        )
        
        import json
        import re
        
        # Clean the response to extract JSON
        response_text = response_text.strip()
        
        # Remove markdown code blocks if present
        if response_text.startswith('```'):
            lines = response_text.split('\n')
            if lines[0].strip().startswith('```'):
                lines = lines[1:]
            if lines and lines[-1].strip() == '```':
                lines = lines[:-1]
            response_text = '\n'.join(lines)
        
        try:
            result = json.loads(response_text)
        except json.JSONDecodeError:
            # Try to find JSON in the text
            json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
            if json_match:
                try:
                    result = json.loads(json_match.group(0))
                except:
                    # Fallback: create basic structure
                    result = {
                        "cards": [{"statement": "Content", "answer": True, "explanation": response_text}]
                    }
            else:
                # If still fails, wrap in basic structure
                result = {
                    "cards": [{"statement": "Content", "answer": True, "explanation": response_text}]
                }
        
        # Ensure cards have required fields
        if "cards" in result:
            for card in result["cards"]:
                if "answer" not in card:
                    card["answer"] = True
                if "explanation" not in card:
                    card["explanation"] = "No explanation provided"
        
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/exam-from-files")
async def exam_from_files(
    request: Request,
    req: ExamRequest,
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_db)
):
    rate_limit(request)
    
    if current_user and not check_quota(db, current_user, "exam"):
        raise HTTPException(status_code=403, detail="Exam generation quota exceeded")
    
    # Require either file_ids or prompt
    if not req.file_ids and not req.prompt:
        raise HTTPException(status_code=400, detail="Please provide either files or a prompt.")
    
    # Fetch file contents from database or in-memory storage
    file_contents = []
    if req.file_ids:
        for file_id in req.file_ids:
            # Try database first
            upload = db.query(Upload).filter(Upload.file_id == file_id).first()
            if upload and upload.content:
                file_contents.append(upload.content)
            # Try in-memory storage for anonymous users
            elif file_id in file_content_store:
                file_contents.append(file_content_store[file_id]["content"])
    
    level_text = get_level_text(req.level)
    
    additional_instructions = ""
    if req.prompt:
        additional_instructions = f"\n\nADDITIONAL USER INSTRUCTIONS:\n{req.prompt}\n\nMake sure to follow these instructions while creating exam questions."
    
    language_instruction = ""
    if req.language == "tr":
        language_instruction = "\n\nIMPORTANT: Generate ALL exam questions in TURKISH language."
    elif req.language == "en":
        language_instruction = "\n\nIMPORTANT: Generate ALL exam questions in ENGLISH language."
    
    system_prompt = """You are a study assistant. Create exam questions intelligently from documents or user topics."""
    
    if file_contents:
        user_prompt = f"""IMPORTANT: First, analyze the document to determine its type:
{language_instruction}

1. If the document contains EXISTING EXAM QUESTIONS/TESTS:
   - Identify the topics, subjects, and difficulty level
   - Generate {req.count} NEW similar questions on the SAME TOPICS
   - DO NOT ask questions about the exam itself (e.g., "What does question 3 ask?")
   - Instead, create new questions about the same subject matter
   - Example: If the exam has SQL database questions, create NEW SQL questions, not "What table is mentioned in the exam?"

2. If the document contains STUDY MATERIAL (notes, textbooks, lectures):
   - Generate {req.count} questions testing knowledge of the content
   - Ask questions about concepts, definitions, and topics in the material

Difficulty level: {level_text}
{additional_instructions}

Follow this EXACT format:

1. Question text here?
A) Option A
B) Option B
C) Option C
D) Option D

2. Question text here?
A) Option A
B) Option B
C) Option C
D) Option D

Cevap Anahtarı:
1-A, 2-B, ...

Generate questions now based on the document type you identified."""
    else:
        # No files, generate from prompt
        user_prompt = f"""Create {req.count} multiple-choice exam questions on the topic requested by the user.
{language_instruction}

USER REQUEST:
{req.prompt}

Difficulty level: {level_text}

Create educational questions that test understanding of the topic.

Follow this EXACT format:

1. Question text here?
A) Option A
B) Option B
C) Option C
D) Option D

2. Question text here?
A) Option A
B) Option B
C) Option C
D) Option D

Cevap Anahtarı:
1-A, 2-B, ..."""
    
    try:
        user_id = current_user.id if current_user else None
        response_text = call_openai_with_context(
            file_contents, 
            f"{system_prompt}\n\n{user_prompt}", 
            temperature=0.0,
            user_id=user_id,
            endpoint="exam",
            db=db
        )
        
        # Check for insufficient context
        if "INSUFFICIENT_CONTEXT" in response_text:
            return {"status": "INSUFFICIENT_CONTEXT", "message": response_text}
        
        # Parse questions and answer key
        questions, answer_key = parse_mcq_questions(response_text)
        
        # If no questions were parsed, return error
        if not questions:
            return {
                "status": "ERROR",
                "message": "Could not generate valid questions from the document. Raw response: " + response_text[:500]
            }
        
        result = {
            "questions": questions,
            "answer_key": answer_key,
            "grounding": [{"number": i+1, "sources": [{"file_id": req.file_ids[0] if req.file_ids else "", "evidence": "Based on uploaded documents"}]} for i in range(len(questions))]
        }
        
        if current_user:
            increment_usage(db, current_user, "exam")
            
            # Save exam
            import json
            exam = Exam(user_id=current_user.id, payload_json=json.dumps(result))
            db.add(exam)
            db.commit()
        
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# ============================================================================
# UNGROUNDED ENDPOINTS
# ============================================================================
@app.post("/ask")
async def ask(
    request: Request,
    req: AskRequest,
    authorization: Optional[str] = Header(None),
    db: Session = Depends(get_db)
):
    rate_limit(request)
    
    # Try to get user if authenticated, otherwise allow limited access
    current_user = None
    if authorization and authorization.startswith("Bearer "):
        token = authorization.split(" ")[1]
        try:
            payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
            user_id_str = payload.get("sub")
            if user_id_str:
                user_id = int(user_id_str)  # Convert string to int
                current_user = db.query(User).filter(User.id == user_id).first()
        except JWTError:
            pass
    
    # Check quota only if user is authenticated
    if current_user and not check_quota(db, current_user, "exam"):
        raise HTTPException(status_code=403, detail="Exam generation quota exceeded. Please upgrade to Premium or try again tomorrow.")
    
    level_text = get_level_text(req.level)
    
    prompt = f"""Create {req.count} multiple-choice questions (A–D) about: {req.prompt}
Difficulty: {level_text}

Follow this EXACT format:

1. Question text here?
A) Option A
B) Option B
C) Option C
D) Option D

2. Question text here?
A) Option A
B) Option B
C) Option C
D) Option D

Cevap Anahtarı:
1-A, 2-B, ..."""
    
    try:
        user_id = current_user.id if current_user else None
        response_text = call_openai_with_context(
            [], 
            prompt, 
            temperature=0.0,
            user_id=user_id,
            endpoint="exam",
            db=db
        )
        
        questions, answer_key = parse_mcq_questions(response_text)
        
        result = {
            "questions": questions,
            "answer_key": answer_key
        }
        
        # Increment usage only if user is authenticated
        if current_user:
            increment_usage(db, current_user, "exam")
        
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/explain")
async def explain(
    request: Request,
    req: ExplainRequest,
    authorization: Optional[str] = Header(None),
    db: Session = Depends(get_db)
):
    rate_limit(request)
    
    # Try to get user if authenticated
    current_user = None
    if authorization and authorization.startswith("Bearer "):
        token = authorization.split(" ")[1]
        try:
            payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
            user_id_str = payload.get("sub")
            if user_id_str:
                user_id = int(user_id_str)  # Convert string to int
                current_user = db.query(User).filter(User.id == user_id).first()
        except JWTError:
            pass
    
    # Check quota only if user is authenticated
    if current_user and not check_quota(db, current_user, "explain"):
        raise HTTPException(status_code=403, detail="Explanation quota exceeded")
    
    # Fetch file contents if file_ids are provided
    file_contents = []
    if req.file_ids:
        for file_id in req.file_ids:
            # Try database first
            upload = db.query(Upload).filter(Upload.file_id == file_id).first()
            if upload and upload.content:
                file_contents.append(upload.content)
            # Try in-memory storage for anonymous users
            elif file_id in file_content_store:
                file_contents.append(file_content_store[file_id]["content"])
    
    prompt = f"""Explain this question concisely:

Question: {req.question}
"""
    
    if req.options:
        for key, val in req.options.items():
            prompt += f"{key}) {val}\n"
    
    if req.selected and req.correct:
        prompt += f"\nSelected: {req.selected}\nCorrect: {req.correct}\n"
    
    if file_contents:
        prompt += """\n\nNOTE: The uploaded document contains reference material. 
If it's an exam/test, use it to understand the topic and context, but explain based on the SUBJECT MATTER, not the document itself.
If it's study material, reference specific information from it.
"""
    
    prompt += "\nProvide a short, targeted explanation in the same language as the question. Justify the correct answer and explain why others are incorrect."
    
    try:
        user_id = current_user.id if current_user else None
        response_text = call_openai_with_context(
            file_contents, 
            prompt, 
            temperature=0.2,
            user_id=user_id,
            endpoint="explain",
            db=db
        )
        
        # Increment usage only if user is authenticated
        if current_user:
            increment_usage(db, current_user, "explain")
        
        return {"explanation": response_text}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/chat")
async def chat(
    request: Request,
    req: ChatRequest,
    authorization: Optional[str] = Header(None),
    db: Session = Depends(get_db)
):
    rate_limit(request)
    
    # Try to get user if authenticated
    current_user = None
    if authorization and authorization.startswith("Bearer "):
        token = authorization.split(" ")[1]
        try:
            payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
            user_id_str = payload.get("sub")
            if user_id_str:
                user_id = int(user_id_str)  # Convert string to int
                current_user = db.query(User).filter(User.id == user_id).first()
        except JWTError:
            pass
    
    # Check quota only if user is authenticated
    if current_user and not check_quota(db, current_user, "chat"):
        raise HTTPException(status_code=403, detail="Chat quota exceeded")
    
    # Fetch file contents if file_ids are provided
    file_contents = []
    if req.file_ids:
        for file_id in req.file_ids:
            # Try database first
            upload = db.query(Upload).filter(Upload.file_id == file_id).first()
            if upload and upload.content:
                file_contents.append(upload.content)
            # Try in-memory storage for anonymous users
            elif file_id in file_content_store:
                file_contents.append(file_content_store[file_id]["content"])
    
    # Build conversation
    conversation = "\n".join([f"{msg.role}: {msg.content}" for msg in req.messages])
    
    try:
        user_id = current_user.id if current_user else None
        response_text = call_openai_with_context(
            file_contents, 
            conversation, 
            temperature=0.7,
            user_id=user_id,
            endpoint="chat",
            db=db
        )
        
        # Increment usage only if user is authenticated
        if current_user:
            increment_usage(db, current_user, "chat")
        
        return {"response": response_text}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# ============================================================================
# BILLING ENDPOINTS (STRIPE)
# ============================================================================
@app.post("/billing/create-checkout-session")
async def create_checkout_session(
    req: CheckoutRequest,
    current_user: User = Depends(get_current_user)
):
    if not STRIPE_SECRET_KEY:
        raise HTTPException(status_code=500, detail="Stripe not configured")
    
    try:
        session = stripe.checkout.Session.create(
            payment_method_types=["card"],
            line_items=[{
                "price": req.priceId,
                "quantity": 1,
            }],
            mode="subscription",
            success_url=req.successUrl,
            cancel_url=req.cancelUrl,
            client_reference_id=str(current_user.id),
        )
        
        return {"sessionId": session.id, "url": session.url}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/billing/webhook")
async def stripe_webhook(request: Request, db: Session = Depends(get_db)):
    if not STRIPE_WEBHOOK_SECRET:
        raise HTTPException(status_code=500, detail="Stripe webhook not configured")
    
    payload = await request.body()
    sig_header = request.headers.get("stripe-signature")
    
    try:
        event = stripe.Webhook.construct_event(
            payload, sig_header, STRIPE_WEBHOOK_SECRET
        )
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))
    
    import json
    
    # Handle checkout.session.completed
    if event["type"] == "checkout.session.completed":
        session = event["data"]["object"]
        user_id = session.get("client_reference_id")
        
        if user_id:
            user = db.query(User).filter(User.id == int(user_id)).first()
            if user:
                user.tier = "premium"
                
                # Get amount from session
                amount_total = session.get("amount_total", 0)  # Amount in cents
                amount = amount_total / 100.0  # Convert to dollars
                
                # Get customer and subscription info
                customer_id = session.get("customer")
                subscription_id = session.get("subscription")
                
                # Store transaction
                transaction = Transaction(
                    user_id=user.id,
                    stripe_session_id=session.get("id"),
                    stripe_customer_id=customer_id,
                    stripe_subscription_id=subscription_id,
                    stripe_payment_intent_id=session.get("payment_intent"),
                    amount=amount,
                    currency=session.get("currency", "usd"),
                    status="completed",
                    tier="premium",
                    event_type="checkout.session.completed",
                    event_metadata=json.dumps(session)
                )
                db.add(transaction)
                db.commit()
    
    # Handle invoice.paid (recurring payments)
    elif event["type"] == "invoice.paid":
        invoice = event["data"]["object"]
        customer_id = invoice.get("customer")
        subscription_id = invoice.get("subscription")
        amount = invoice.get("amount_paid", 0) / 100.0  # Convert from cents to dollars
        
        # Find user by customer_id from previous transactions
        transaction = db.query(Transaction).filter(
            Transaction.stripe_customer_id == customer_id
        ).first()
        
        if transaction:
            user_id = transaction.user_id
            # Store recurring payment transaction
            recurring_transaction = Transaction(
                user_id=user_id,
                stripe_session_id=invoice.get("id"),  # Use invoice ID as session ID
                stripe_customer_id=customer_id,
                stripe_subscription_id=subscription_id,
                stripe_payment_intent_id=invoice.get("payment_intent"),
                amount=amount,
                currency=invoice.get("currency", "usd"),
                status="completed",
                tier="premium",
                event_type="invoice.paid",
                event_metadata=json.dumps(invoice)
            )
            db.add(recurring_transaction)
            db.commit()
    
    # Handle payment failures
    elif event["type"] == "invoice.payment_failed":
        invoice = event["data"]["object"]
        customer_id = invoice.get("customer")
        amount = invoice.get("amount_due", 0) / 100.0
        
        transaction = db.query(Transaction).filter(
            Transaction.stripe_customer_id == customer_id
        ).first()
        
        if transaction:
            failed_transaction = Transaction(
                user_id=transaction.user_id,
                stripe_session_id=invoice.get("id"),
                stripe_customer_id=customer_id,
                stripe_subscription_id=invoice.get("subscription"),
                amount=amount,
                currency=invoice.get("currency", "usd"),
                status="failed",
                tier="premium",
                event_type="invoice.payment_failed",
                event_metadata=json.dumps(invoice)
            )
            db.add(failed_transaction)
            db.commit()
    
    return {"status": "success"}

# ============================================================================
# HISTORY ENDPOINTS
# ============================================================================
@app.post("/history", response_model=HistoryItemResponse)
async def save_history(
    item: HistoryItemCreate,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Save a history item for the authenticated user"""
    import json
    
    # Validate folder_id if provided
    if item.folder_id:
        folder = db.query(Folder).filter(
            Folder.id == item.folder_id,
            Folder.user_id == current_user.id
        ).first()
        if not folder:
            raise HTTPException(status_code=404, detail="Folder not found")
    
    history_entry = History(
        user_id=current_user.id,
        type=item.type,
        title=item.title,
        data_json=json.dumps(item.data),
        score_json=json.dumps(item.score) if item.score else None,
        folder_id=item.folder_id
    )
    db.add(history_entry)
    db.commit()
    db.refresh(history_entry)
    
    return {
        "id": history_entry.id,
        "type": history_entry.type,
        "title": history_entry.title,
        "data": json.loads(history_entry.data_json),
        "score": json.loads(history_entry.score_json) if history_entry.score_json else None,
        "folder_id": history_entry.folder_id,
        "timestamp": int(history_entry.created_at.timestamp() * 1000)
    }

@app.get("/history", response_model=List[HistoryItemResponse])
async def get_history(
    folder_id: Optional[int] = None,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get history items for the authenticated user, optionally filtered by folder"""
    import json
    
    query = db.query(History).filter(History.user_id == current_user.id)
    
    # Filter by folder if provided
    if folder_id is not None:
        if folder_id == 0:
            # Special case: folder_id=0 means "no folder" (uncategorized)
            query = query.filter(History.folder_id == None)
        else:
            query = query.filter(History.folder_id == folder_id)
    
    history_entries = query.order_by(History.created_at.desc()).all()
    
    return [
        {
            "id": entry.id,
            "type": entry.type,
            "title": entry.title,
            "data": json.loads(entry.data_json),
            "score": json.loads(entry.score_json) if entry.score_json else None,
            "folder_id": entry.folder_id,
            "timestamp": int(entry.created_at.timestamp() * 1000)
        }
        for entry in history_entries
    ]

@app.put("/history/{history_id}", response_model=HistoryItemResponse)
async def update_history_item(
    history_id: int,
    update: HistoryItemUpdate,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Update a specific history item (e.g., add score after completion or move to folder)"""
    import json
    
    history_entry = db.query(History).filter(
        History.id == history_id,
        History.user_id == current_user.id
    ).first()
    
    if not history_entry:
        raise HTTPException(status_code=404, detail="History item not found")
    
    # Validate folder_id if provided
    if update.folder_id is not None:
        if update.folder_id > 0:  # 0 or None means remove from folder
            folder = db.query(Folder).filter(
                Folder.id == update.folder_id,
                Folder.user_id == current_user.id
            ).first()
            if not folder:
                raise HTTPException(status_code=404, detail="Folder not found")
    
    # Update fields if provided
    if update.title is not None:
        history_entry.title = update.title
    if update.data is not None:
        history_entry.data_json = json.dumps(update.data)
    if update.score is not None:
        history_entry.score_json = json.dumps(update.score)
    if update.folder_id is not None:
        history_entry.folder_id = update.folder_id if update.folder_id > 0 else None
    
    db.commit()
    db.refresh(history_entry)
    
    return {
        "id": history_entry.id,
        "type": history_entry.type,
        "title": history_entry.title,
        "data": json.loads(history_entry.data_json),
        "score": json.loads(history_entry.score_json) if history_entry.score_json else None,
        "folder_id": history_entry.folder_id,
        "timestamp": int(history_entry.created_at.timestamp() * 1000)
    }

@app.delete("/history/{history_id}")
async def delete_history_item(
    history_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Delete a specific history item"""
    history_entry = db.query(History).filter(
        History.id == history_id,
        History.user_id == current_user.id
    ).first()
    
    if not history_entry:
        raise HTTPException(status_code=404, detail="History item not found")
    
    db.delete(history_entry)
    db.commit()
    
    return {"status": "success"}

@app.delete("/history")
async def clear_history(
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Clear all history for the authenticated user"""
    db.query(History).filter(History.user_id == current_user.id).delete()
    db.commit()
    
    return {"status": "success"}

# ============================================================================
# FOLDER ENDPOINTS
# ============================================================================
@app.post("/folders", response_model=FolderResponse)
async def create_folder(
    folder: FolderCreate,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Create a new folder for organizing history items"""
    new_folder = Folder(
        user_id=current_user.id,
        name=folder.name,
        color=folder.color,
        icon=folder.icon
    )
    db.add(new_folder)
    db.commit()
    db.refresh(new_folder)
    
    # Count items in this folder
    item_count = db.query(History).filter(
        History.user_id == current_user.id,
        History.folder_id == new_folder.id
    ).count()
    
    return {
        "id": new_folder.id,
        "name": new_folder.name,
        "color": new_folder.color,
        "icon": new_folder.icon,
        "created_at": new_folder.created_at.isoformat(),
        "item_count": item_count
    }

@app.get("/folders", response_model=List[FolderResponse])
async def get_folders(
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get all folders for the authenticated user"""
    folders = db.query(Folder).filter(
        Folder.user_id == current_user.id
    ).order_by(Folder.created_at.desc()).all()
    
    result = []
    for folder in folders:
        # Count items in this folder
        item_count = db.query(History).filter(
            History.user_id == current_user.id,
            History.folder_id == folder.id
        ).count()
        
        result.append({
            "id": folder.id,
            "name": folder.name,
            "color": folder.color,
            "icon": folder.icon,
            "created_at": folder.created_at.isoformat(),
            "item_count": item_count
        })
    
    return result

@app.put("/folders/{folder_id}", response_model=FolderResponse)
async def update_folder(
    folder_id: int,
    update: FolderUpdate,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Update a folder"""
    folder = db.query(Folder).filter(
        Folder.id == folder_id,
        Folder.user_id == current_user.id
    ).first()
    
    if not folder:
        raise HTTPException(status_code=404, detail="Folder not found")
    
    # Update fields if provided
    if update.name is not None:
        folder.name = update.name
    if update.color is not None:
        folder.color = update.color
    if update.icon is not None:
        folder.icon = update.icon
    
    db.commit()
    db.refresh(folder)
    
    # Count items in this folder
    item_count = db.query(History).filter(
        History.user_id == current_user.id,
        History.folder_id == folder.id
    ).count()
    
    return {
        "id": folder.id,
        "name": folder.name,
        "color": folder.color,
        "icon": folder.icon,
        "created_at": folder.created_at.isoformat(),
        "item_count": item_count
    }

@app.delete("/folders/{folder_id}")
async def delete_folder(
    folder_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Delete a folder (history items will become uncategorized)"""
    folder = db.query(Folder).filter(
        Folder.id == folder_id,
        Folder.user_id == current_user.id
    ).first()
    
    if not folder:
        raise HTTPException(status_code=404, detail="Folder not found")
    
    # Remove folder_id from all history items in this folder
    db.query(History).filter(
        History.user_id == current_user.id,
        History.folder_id == folder_id
    ).update({"folder_id": None})
    
    db.delete(folder)
    db.commit()
    
    return {"status": "success"}

# ============================================================================
# ADMIN ENDPOINTS
# ============================================================================
@app.get("/admin/users", response_model=List[UserResponse])
async def get_all_users(
    skip: int = 0,
    limit: int = 100,
    admin_user: User = Depends(get_admin_user),
    db: Session = Depends(get_db)
):
    """Get all users (admin only)"""
    users = db.query(User).offset(skip).limit(limit).all()
    return [
        {
            "id": user.id,
            "email": user.email,
            "name": user.name,
            "surname": user.surname,
            "tier": user.tier,
            "is_admin": bool(user.is_admin) if user.is_admin else False,
            "created_at": user.created_at.isoformat()
        }
        for user in users
    ]

@app.get("/admin/users/{user_id}")
async def get_user(
    user_id: int,
    admin_user: User = Depends(get_admin_user),
    db: Session = Depends(get_db)
):
    """Get a specific user by ID (admin only)"""
    user = db.query(User).filter(User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Get usage statistics
    today = date.today()
    usage_data = {}
    for kind in ["exam", "explain", "chat", "upload"]:
        usage = db.query(Usage).filter(
            Usage.user_id == user.id,
            Usage.kind == kind,
            Usage.date == today
        ).first()
        usage_data[kind] = usage.count if usage else 0
    
    # Get total history items
    history_count = db.query(History).filter(History.user_id == user.id).count()
    
    # Get total uploads
    uploads_count = db.query(Upload).filter(Upload.user_id == user.id).count()
    
    return {
        "id": user.id,
        "email": user.email,
        "name": user.name,
        "surname": user.surname,
        "tier": user.tier,
        "is_admin": bool(user.is_admin) if user.is_admin else False,
        "created_at": user.created_at.isoformat(),
        "usage": usage_data,
        "history_count": history_count,
        "uploads_count": uploads_count
    }

@app.put("/admin/users/{user_id}", response_model=UserResponse)
async def update_user(
    user_id: int,
    user_update: UserUpdate,
    admin_user: User = Depends(get_admin_user),
    db: Session = Depends(get_db)
):
    """Update a user (admin only)"""
    user = db.query(User).filter(User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Prevent admin from removing their own admin status
    if user_id == admin_user.id and user_update.is_admin is False:
        raise HTTPException(status_code=400, detail="Cannot remove your own admin status")
    
    if user_update.name is not None:
        user.name = user_update.name
    if user_update.surname is not None:
        user.surname = user_update.surname
    if user_update.tier is not None:
        if user_update.tier not in ["free", "premium", "standard", "pro"]:
            raise HTTPException(status_code=400, detail="Invalid tier. Must be 'free', 'standard', 'premium', or 'pro'")
        user.tier = user_update.tier
    if user_update.is_admin is not None:
        user.is_admin = 1 if user_update.is_admin else 0
    
    db.commit()
    db.refresh(user)
    
    return {
        "id": user.id,
        "email": user.email,
        "name": user.name,
        "surname": user.surname,
        "tier": user.tier,
        "is_admin": bool(user.is_admin) if user.is_admin else False,
        "created_at": user.created_at.isoformat()
    }

@app.delete("/admin/users/{user_id}")
async def delete_user(
    user_id: int,
    admin_user: User = Depends(get_admin_user),
    db: Session = Depends(get_db)
):
    """Delete a user (admin only)"""
    user = db.query(User).filter(User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Prevent admin from deleting themselves
    if user_id == admin_user.id:
        raise HTTPException(status_code=400, detail="Cannot delete your own account")
    
    # Delete related data
    db.query(History).filter(History.user_id == user_id).delete()
    db.query(Upload).filter(Upload.user_id == user_id).delete()
    db.query(Usage).filter(Usage.user_id == user_id).delete()
    db.query(Exam).filter(Exam.user_id == user_id).delete()
    db.query(Folder).filter(Folder.user_id == user_id).delete()
    
    db.delete(user)
    db.commit()
    
    return {"status": "success", "message": f"User {user_id} deleted successfully"}

@app.get("/admin/stats")
async def get_admin_stats(
    admin_user: User = Depends(get_admin_user),
    db: Session = Depends(get_db)
):
    """Get admin statistics (admin only)"""
    # Total users
    total_users = db.query(User).count()
    free_users = db.query(User).filter(User.tier == "free").count()
    premium_users = db.query(User).filter(User.tier.in_(["premium", "pro"])).count()
    admin_users = db.query(User).filter(User.is_admin == 1).count()
    
    # Usage statistics
    today = date.today()
    today_usage = db.query(Usage).filter(Usage.date == today).all()
    usage_by_kind = {}
    for usage in today_usage:
        if usage.kind not in usage_by_kind:
            usage_by_kind[usage.kind] = 0
        usage_by_kind[usage.kind] += usage.count
    
    # Total history items
    total_history = db.query(History).count()
    
    # Total uploads
    total_uploads = db.query(Upload).count()
    
    # Recent registrations (last 7 days)
    seven_days_ago = datetime.utcnow() - timedelta(days=7)
    recent_registrations = db.query(User).filter(User.created_at >= seven_days_ago).count()
    
    return {
        "users": {
            "total": total_users,
            "free": free_users,
            "premium": premium_users,
            "admin": admin_users
        },
        "usage_today": usage_by_kind,
        "content": {
            "total_history": total_history,
            "total_uploads": total_uploads
        },
        "recent_registrations_7d": recent_registrations
    }

@app.get("/admin/quality-stats")
async def get_quality_stats_endpoint(
    days: int = 7,
    admin_user: User = Depends(get_admin_user),
    db: Session = Depends(get_db)
):
    """Get quality statistics for recent summaries (admin only)"""
    try:
        from app.services.telemetry import get_quality_stats
        stats = get_quality_stats(db, days=days)
        return stats
    except Exception as e:
        return {"error": str(e)}


@app.get("/admin/low-quality-patterns")
async def get_low_quality_patterns_endpoint(
    threshold: float = 0.6,
    limit: int = 10,
    admin_user: User = Depends(get_admin_user),
    db: Session = Depends(get_db)
):
    """Get patterns in low-quality summaries for improvement (admin only)"""
    try:
        from app.services.telemetry import get_low_quality_patterns
        patterns = get_low_quality_patterns(db, threshold=threshold, limit=limit)
        return {"patterns": patterns, "threshold": threshold}
    except Exception as e:
        return {"error": str(e)}


@app.get("/admin/check-user/{email}")
async def check_user_admin_status(
    email: str,
    db: Session = Depends(get_db)
):
    """Check if a user exists and their admin status (public endpoint for debugging)"""
    try:
        user = db.query(User).filter(User.email == email).first()
        if not user:
            return {
                "exists": False,
                "message": f"User with email '{email}' not found"
            }
        
        is_admin_value = getattr(user, 'is_admin', 0) or 0
        is_admin = bool(int(is_admin_value)) if is_admin_value else False
        
        return {
            "exists": True,
            "user": {
                "id": user.id,
                "email": user.email,
                "name": getattr(user, 'name', None),
                "surname": getattr(user, 'surname', None),
                "tier": getattr(user, 'tier', 'free'),
                "is_admin": is_admin,
                "is_admin_raw": is_admin_value
            }
        }
    except Exception as e:
        return {
            "exists": False,
            "error": str(e)
        }

@app.post("/admin/bootstrap-admin")
async def bootstrap_admin(
    request: BootstrapAdminRequest,
    db: Session = Depends(get_db)
):
    """
    Bootstrap the first admin user. 
    Only works if:
    1. No admins exist yet, OR
    2. A valid secret key is provided (from BOOTSTRAP_SECRET_KEY env var)
    """
    try:
        # Check if any admins exist (is_admin is integer: 0 or 1)
        admin_count = db.query(User).filter(User.is_admin == 1).count()
        
        print(f"[BOOTSTRAP] Admin count: {admin_count}")
        print(f"[BOOTSTRAP] Request email: {request.email}")
        print(f"[BOOTSTRAP] Secret key provided: {bool(request.secret_key)}")
        
        # If admins exist, require secret key
        if admin_count > 0:
            expected_secret = os.getenv("BOOTSTRAP_SECRET_KEY")
            if not expected_secret:
                raise HTTPException(
                    status_code=403, 
                    detail="Admins already exist. Set BOOTSTRAP_SECRET_KEY environment variable to bootstrap additional admins."
                )
            
            if not request.secret_key or request.secret_key != expected_secret:
                raise HTTPException(
                    status_code=403,
                    detail="Invalid secret key. Cannot bootstrap admin when admins already exist."
                )
        
        # Find the user
        user = db.query(User).filter(User.email == request.email).first()
        if not user:
            raise HTTPException(status_code=404, detail=f"User with email '{request.email}' not found")
        
        print(f"[BOOTSTRAP] User found: ID={user.id}, Current is_admin={getattr(user, 'is_admin', 'N/A')}")
        
        # Make user admin
        user.is_admin = 1
        db.commit()
        db.refresh(user)
        
        # Verify the change
        is_admin_after = getattr(user, 'is_admin', 0) or 0
        print(f"[BOOTSTRAP] After update: is_admin={is_admin_after}")
        
        return {
            "status": "success",
            "message": f"User '{request.email}' is now an admin",
            "user": {
                "id": user.id,
                "email": user.email,
                "name": user.name,
                "surname": user.surname,
                "is_admin": bool(int(is_admin_after)) if is_admin_after else False,
                "is_admin_raw": is_admin_after
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error bootstrapping admin: {str(e)}")

@app.get("/admin/migrate-database")
@app.post("/admin/migrate-database")
async def migrate_database():
    """Manually trigger database migration"""
    try:
        run_migration()
        return {"status": "success", "message": "Database migrated successfully - all columns and tables added"}
    except Exception as e:
        error_msg = str(e)
        return {"status": "error", "message": error_msg}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)

@app.get("/admin/clear-cache")
@app.delete("/admin/clear-cache")
async def clear_cache(
    admin_user: User = Depends(get_admin_user),
    db: Session = Depends(get_db)
):
    """Clear all summary cache (admin only)"""
    try:
        from app.services.cache import SummaryCache
        deleted = db.query(SummaryCache).delete()
        db.commit()
        return {"status": "success", "deleted": deleted, "message": "Cache cleared successfully"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/admin/transactions")
async def get_transactions(
    skip: int = 0,
    limit: int = 100,
    user_id: Optional[int] = None,
    admin_user: User = Depends(get_admin_user),
    db: Session = Depends(get_db)
):
    """Get all transactions (admin only)"""
    query = db.query(Transaction)
    if user_id:
        query = query.filter(Transaction.user_id == user_id)
    transactions = query.order_by(Transaction.created_at.desc()).offset(skip).limit(limit).all()
    
    return [
        {
            "id": t.id,
            "user_id": t.user_id,
            "user_email": db.query(User).filter(User.id == t.user_id).first().email if db.query(User).filter(User.id == t.user_id).first() else None,
            "stripe_session_id": t.stripe_session_id,
            "stripe_customer_id": t.stripe_customer_id,
            "stripe_subscription_id": t.stripe_subscription_id,
            "amount": t.amount,
            "currency": t.currency,
            "status": t.status,
            "tier": t.tier,
            "event_type": t.event_type,
            "created_at": t.created_at.isoformat()
        }
        for t in transactions
    ]

@app.get("/admin/token-usage")
async def get_token_usage(
    skip: int = 0,
    limit: int = 100,
    user_id: Optional[int] = None,
    endpoint: Optional[str] = None,
    admin_user: User = Depends(get_admin_user),
    db: Session = Depends(get_db)
):
    """Get token usage statistics (admin only)"""
    query = db.query(TokenUsage)
    if user_id:
        query = query.filter(TokenUsage.user_id == user_id)
    if endpoint:
        query = query.filter(TokenUsage.endpoint == endpoint)
    usage_records = query.order_by(TokenUsage.created_at.desc()).offset(skip).limit(limit).all()
    
    return [
        {
            "id": u.id,
            "user_id": u.user_id,
            "user_email": db.query(User).filter(User.id == u.user_id).first().email if u.user_id and db.query(User).filter(User.id == u.user_id).first() else None,
            "endpoint": u.endpoint,
            "model": u.model,
            "input_tokens": u.input_tokens,
            "output_tokens": u.output_tokens,
            "total_tokens": u.total_tokens,
            "estimated_cost": u.estimated_cost,
            "created_at": u.created_at.isoformat()
        }
        for u in usage_records
    ]

@app.get("/admin/revenue")
async def get_revenue_stats(
    days: int = 30,
    admin_user: User = Depends(get_admin_user),
    db: Session = Depends(get_db)
):
    """Get revenue statistics (admin only)"""
    from datetime import datetime, timedelta
    cutoff_date = datetime.utcnow() - timedelta(days=days)
    
    # Total revenue
    total_revenue = db.query(Transaction).filter(
        Transaction.status == "completed",
        Transaction.created_at >= cutoff_date
    ).with_entities(func.sum(Transaction.amount)).scalar() or 0.0
    
    # Revenue by period (daily for last 30 days)
    daily_revenue = []
    for i in range(days):
        date = datetime.utcnow() - timedelta(days=i)
        day_start = date.replace(hour=0, minute=0, second=0, microsecond=0)
        day_end = day_start + timedelta(days=1)
        
        revenue = db.query(Transaction).filter(
            Transaction.status == "completed",
            Transaction.created_at >= day_start,
            Transaction.created_at < day_end
        ).with_entities(func.sum(Transaction.amount)).scalar() or 0.0
        
        daily_revenue.append({
            "date": day_start.date().isoformat(),
            "revenue": revenue
        })
    
    daily_revenue.reverse()  # Oldest to newest
    
    # Revenue by tier
    revenue_by_tier = {}
    tiers = db.query(Transaction.tier).distinct().all()
    for (tier,) in tiers:
        revenue = db.query(Transaction).filter(
            Transaction.status == "completed",
            Transaction.tier == tier,
            Transaction.created_at >= cutoff_date
        ).with_entities(func.sum(Transaction.amount)).scalar() or 0.0
        revenue_by_tier[tier] = revenue
    
    # Total transactions
    total_transactions = db.query(Transaction).filter(
        Transaction.status == "completed",
        Transaction.created_at >= cutoff_date
    ).count()
    
    # Failed transactions
    failed_transactions = db.query(Transaction).filter(
        Transaction.status == "failed",
        Transaction.created_at >= cutoff_date
    ).count()
    
    # Total token costs
    total_token_cost = db.query(TokenUsage).filter(
        TokenUsage.created_at >= cutoff_date
    ).with_entities(func.sum(TokenUsage.estimated_cost)).scalar() or 0.0
    
    # Token usage by endpoint
    token_usage_by_endpoint = {}
    endpoints = db.query(TokenUsage.endpoint).distinct().all()
    for (endpoint,) in endpoints:
        usage = db.query(TokenUsage).filter(
            TokenUsage.endpoint == endpoint,
            TokenUsage.created_at >= cutoff_date
        ).with_entities(
            func.sum(TokenUsage.total_tokens).label("total_tokens"),
            func.sum(TokenUsage.estimated_cost).label("total_cost")
        ).first()
        if usage:
            token_usage_by_endpoint[endpoint] = {
                "total_tokens": usage.total_tokens or 0,
                "total_cost": usage.total_cost or 0.0
            }
    
    # Token usage by user (top 10)
    top_users = db.query(
        TokenUsage.user_id,
        func.sum(TokenUsage.total_tokens).label("total_tokens"),
        func.sum(TokenUsage.estimated_cost).label("total_cost")
    ).filter(
        TokenUsage.created_at >= cutoff_date,
        TokenUsage.user_id.isnot(None)
    ).group_by(TokenUsage.user_id).order_by(func.sum(TokenUsage.estimated_cost).desc()).limit(10).all()
    
    top_users_list = []
    for user_id, total_tokens, total_cost in top_users:
        user = db.query(User).filter(User.id == user_id).first()
        top_users_list.append({
            "user_id": user_id,
            "user_email": user.email if user else None,
            "total_tokens": total_tokens or 0,
            "total_cost": total_cost or 0.0
        })
    
    return {
        "period_days": days,
        "total_revenue": total_revenue,
        "total_transactions": total_transactions,
        "failed_transactions": failed_transactions,
        "daily_revenue": daily_revenue,
        "revenue_by_tier": revenue_by_tier,
        "total_token_cost": total_token_cost,
        "token_usage_by_endpoint": token_usage_by_endpoint,
        "top_users_by_token_cost": top_users_list,
        "net_revenue": total_revenue - total_token_cost
    }

@app.get("/admin/users/{user_id}/details")
async def get_user_detailed_info(
    user_id: int,
    admin_user: User = Depends(get_admin_user),
    db: Session = Depends(get_db)
):
    """Get detailed user information including transactions and token usage (admin only)"""
    user = db.query(User).filter(User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Get user transactions
    transactions = db.query(Transaction).filter(Transaction.user_id == user_id).order_by(Transaction.created_at.desc()).all()
    total_paid = db.query(Transaction).filter(
        Transaction.user_id == user_id,
        Transaction.status == "completed"
    ).with_entities(func.sum(Transaction.amount)).scalar() or 0.0
    
    # Get token usage
    token_usage = db.query(TokenUsage).filter(TokenUsage.user_id == user_id).all()
    total_tokens = db.query(TokenUsage).filter(TokenUsage.user_id == user_id).with_entities(func.sum(TokenUsage.total_tokens)).scalar() or 0
    total_token_cost = db.query(TokenUsage).filter(TokenUsage.user_id == user_id).with_entities(func.sum(TokenUsage.estimated_cost)).scalar() or 0.0
    
    # Token usage by endpoint
    token_by_endpoint = {}
    endpoints = db.query(TokenUsage.endpoint).filter(TokenUsage.user_id == user_id).distinct().all()
    for (endpoint,) in endpoints:
        usage = db.query(TokenUsage).filter(
            TokenUsage.user_id == user_id,
            TokenUsage.endpoint == endpoint
        ).with_entities(
            func.sum(TokenUsage.total_tokens).label("total_tokens"),
            func.sum(TokenUsage.estimated_cost).label("total_cost"),
            func.count().label("count")
        ).first()
        if usage:
            token_by_endpoint[endpoint] = {
                "total_tokens": usage.total_tokens or 0,
                "total_cost": usage.total_cost or 0.0,
                "request_count": usage.count or 0
            }
    
    # Usage statistics
    today = date.today()
    usage_data = {}
    for kind in ["exam", "explain", "chat", "upload"]:
        usage = db.query(Usage).filter(
            Usage.user_id == user.id,
            Usage.kind == kind,
            Usage.date == today
        ).first()
        usage_data[kind] = usage.count if usage else 0
    
    # Total history items
    history_count = db.query(History).filter(History.user_id == user.id).count()
    
    # Total uploads
    uploads_count = db.query(Upload).filter(Upload.user_id == user.id).count()
    
    return {
        "user": {
            "id": user.id,
            "email": user.email,
            "name": user.name,
            "surname": user.surname,
            "tier": user.tier,
            "is_admin": bool(user.is_admin) if user.is_admin else False,
            "created_at": user.created_at.isoformat()
        },
        "transactions": [
            {
                "id": t.id,
                "amount": t.amount,
                "currency": t.currency,
                "status": t.status,
                "tier": t.tier,
                "event_type": t.event_type,
                "created_at": t.created_at.isoformat()
            }
            for t in transactions
        ],
        "total_paid": total_paid,
        "token_usage": {
            "total_tokens": total_tokens,
            "total_cost": total_token_cost,
            "by_endpoint": token_by_endpoint
        },
        "usage_today": usage_data,
        "history_count": history_count,
        "uploads_count": uploads_count
    }
